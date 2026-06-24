import asyncio
import base64
import json
import mimetypes
import os
import random
import re
import subprocess
import tempfile
import time
import xml.etree.ElementTree as ET
from collections import deque
from datetime import datetime, timezone
from html import unescape
from io import BytesIO
from pathlib import Path
from statistics import mean, pstdev
from typing import Any
from uuid import uuid4

import httpx
from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware

try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover - handled gracefully at runtime
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - handled gracefully at runtime
    Presentation = None

try:
    import fitz
except ImportError:  # pragma: no cover - handled gracefully at runtime
    fitz = None

from .config import (
    GEMINI_STATUS_CACHE_SECONDS,
    MAX_PRESENTATION_UPLOAD_BYTES,
    MAX_PRESENTATION_UPLOAD_FILES,
    MAX_VISION_MATERIAL_IMAGES,
    MAX_VISION_PAGES_PER_FILE,
    MAX_SCRIPT_FILE_BYTES,
    OPENAI_AUDIENCE_MAX_SESSION_CALLS,
    OPENAI_AUDIENCE_MIN_INTERVAL_SECONDS,
    OPENAI_AUDIENCE_MODEL,
    TEXT_EXTENSIONS,
    SUPPORTED_SCRIPT_EXTENSIONS,
    get_gemini_model,
)
from .models import AudienceChatRequest, FinishSessionRequest, GeminiRateLimitError, ImportedScriptResponse, MetricSample, SessionState
from .runtime_state import (
    ai_status_cache,
    audience_agent_sessions,
    gemini_limiter,
    gemini_runtime_state,
    openai_audience_limiter,
    openai_runtime_state,
    sessions,
)
from .session_store import append_metric as persist_metric
from .session_store import init_session_store, load_session as load_persisted_session, save_session
from .services.gemini_service import call_gemini_api
from .services.reference_compare_service import build_user_profile_from_samples, compare_with_reference
from .services.reference_service import build_reference_comparison, build_reference_video
from .services.text_service import (
    clamp,
    count_syllables,
    extract_script_text,
    file_extension,
    format_seconds,
    normalize_imported_text,
    recent_excerpt,
    score_distance,
    script_quality,
    tokenize,
    transcript_delta,
)


app = FastAPI(title="Presentation Practice API")

cors_origins = os.getenv(
    "BACKEND_CORS_ORIGINS",
    "http://localhost:5173,http://127.0.0.1:5173",
).split(",")
cors_origin_regex = os.getenv(
    "BACKEND_CORS_ORIGIN_REGEX",
    r"http://(localhost|127\.0\.0\.1):\d+",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=[origin.strip() for origin in cors_origins if origin.strip()],
    allow_origin_regex=cors_origin_regex,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

init_session_store()


def get_session_state(session_id: str) -> SessionState | None:
    session = sessions.get(session_id)
    if session:
        return session
    restored_session = load_persisted_session(session_id)
    if restored_session:
        sessions[session_id] = restored_session
    return restored_session


AUDIENCE_AGENTS = {
    "민서": {
        "role": "공감형 청중",
        "reaction": "attentive",
        "prompt": (
            "너는 공감형 청중 민서다. 발표자의 긴장을 낮춰주는 따뜻한 말투를 쓴다. "
            "칭찬을 먼저 짧게 하고, 필요하면 바로 적용할 수 있는 작은 조언 하나만 덧붙인다."
        ),
        "fallback": {
            "opening": "좋아요. 첫 문장만 천천히 잡고 들어가면 더 안정적이에요.",
            "focused": "지금 흐름 좋아요. 이 속도 그대로 핵심어만 살려주세요.",
            "impressed": "방금 전달이 잘 됐어요. 그 톤으로 이어가면 좋아요.",
            "longSilence": "잠깐 멈췄지만 괜찮아요. 다음 핵심 문장으로 이어가면 돼요.",
            "default": "전달은 잡혀 있어요. 다음 문장 첫 단어만 더 또렷하게 시작해보세요.",
        },
    },
    "준": {
        "role": "분석형 청중",
        "reaction": "tooFast",
        "prompt": (
            "너는 분석형 청중 준이다. 말투는 차분하고 구체적이다. "
            "속도, 침묵, 핵심어 중 하나를 근거로 짧은 피드백을 준다. 감탄보다 수정 포인트가 중요하다."
        ),
        "fallback": {
            "tooFast": "속도가 올라갔어요. 핵심어 앞에서 반 박자만 늦춰주세요.",
            "unclear": "문장 끝이 살짝 흐립니다. 끝 음절을 조금만 더 닫아주세요.",
            "offScript": "대본 핵심어가 약해졌어요. 주제어를 한 번 다시 꺼내주세요.",
            "default": "지표는 괜찮습니다. 한 문장 안의 속도만 더 일정하게 가보세요.",
        },
    },
    "하린": {
        "role": "표현형 청중",
        "reaction": "excited",
        "prompt": (
            "너는 표현형 청중 하린이다. 생동감 있고 몰입감 있는 반응을 한다. "
            "청중 입장에서 어디가 잘 들렸는지, 어디를 더 강조하면 좋은지 감각적으로 말한다."
        ),
        "fallback": {
            "focused": "지금 리듬 좋아요. 중요한 단어에 힘을 조금 더 주세요.",
            "impressed": "방금 구간은 잘 들어왔어요. 그 포인트를 한 번 더 밀어주세요.",
            "offScript": "이야기 방향이 살짝 퍼졌어요. 제목 키워드로 다시 모아보면 좋아요.",
            "tooSlow": "리듬이 조금 처졌어요. 다음 문장은 더 밝게 밀고 가봐요.",
            "default": "톤은 좋아요. 포인트 문장 하나만 더 선명하게 들려주세요.",
        },
    },
    "도윤": {
        "role": "차분형 청중",
        "reaction": "tooSlow",
        "prompt": (
            "너는 차분형 청중 도윤이다. 과장하지 않고 안정적인 말투를 쓴다. "
            "호흡, 쉬는 타이밍, 복구 문장처럼 발표 흐름을 정돈하는 피드백을 준다."
        ),
        "fallback": {
            "tooSlow": "호흡이 길어졌어요. 다음 문장으로 조금 더 빨리 넘어가도 좋습니다.",
            "longSilence": "침묵이 길었습니다. 준비한 연결 문장으로 바로 회복해보세요.",
            "unclear": "소리는 들리지만 문장 경계가 약해요. 짧게 끊어 말해보세요.",
            "default": "전체 흐름은 차분합니다. 쉬는 위치만 조금 더 의도적으로 잡아보세요.",
        },
    },
}

SITUATION_AGENT_POOL = {
    "opening": ["민서", "도윤"],
    "focused": ["도윤", "민서"],
    "impressed": ["하린", "민서"],
    "tooFast": ["준", "도윤"],
    "tooSlow": ["도윤", "하린"],
    "longSilence": ["도윤", "민서"],
    "unclear": ["준", "민서"],
    "offScript": ["준", "하린"],
}

SITUATION_REACTIONS = {
    "opening": "attentive",
    "focused": "attentive",
    "impressed": "excited",
    "tooFast": "tooFast",
    "tooSlow": "sleepy",
    "longSilence": "sleepy",
    "unclear": "confused",
    "offScript": "confused",
}

AUDIENCE_STATE_CRITERIA = {
    "opening": "평상: 발표 시작 직후이거나 아직 충분한 발화 데이터가 쌓이지 않은 상태입니다.",
    "focused": "집중: 속도와 침묵이 안정적이고 대본 핵심어가 어느 정도 유지되는 상태입니다.",
    "impressed": "감탄: 속도가 적절하고 핵심어 반영도가 높아 청중이 몰입하기 좋은 상태입니다.",
    "tooFast": "의문: 말 속도가 빨라 핵심어가 지나가거나 청중이 따라가기 어려운 상태입니다.",
    "tooSlow": "졸림: 느린 진행이나 짧은 침묵 누적으로 청중 집중도가 내려가는 상태입니다.",
    "longSilence": "졸림: 긴 침묵이 이어져 청중 집중도가 크게 내려가는 상태입니다.",
    "unclear": "의문: 목소리는 감지되지만 음성 인식 결과가 충분히 따라오지 않는 상태입니다.",
    "offScript": "의문: 현재 발화가 대본의 핵심어와 멀어져 메시지가 흐려지는 상태입니다.",
}

AUDIENCE_AGENT_STATE_CRITERIA = {
    "민서": "공감형 청중입니다. 짧은 실수에는 바로 부정적으로 반응하지 않고, 긴 침묵이나 명확한 이탈이 이어질 때만 걱정합니다.",
    "준": "분석형 청중입니다. 빠른 속도, 불명확한 발음, 대본 핵심어 이탈에 가장 민감하게 반응합니다.",
    "하린": "표현형 청중입니다. 핵심어가 잘 들어오고 리듬이 살아날 때 가장 빠르게 감탄합니다.",
    "도윤": "차분형 청중입니다. 침묵, 느린 진행, 호흡의 늘어짐에 가장 민감하게 반응합니다.",
}

SITUATION_FALLBACK_TEXT = {
    "opening": "좋아요. 첫 문장만 천천히 잡고 들어가면 더 안정적이에요.",
    "focused": "지금 흐름이 안정적이에요. 이 리듬을 유지해 주세요.",
    "impressed": "방금 포인트가 잘 들어왔어요. 중요한 단어에 힘을 더 주세요.",
    "tooFast": "속도가 올라갔어요. 핵심어 앞에서 반 박자만 늦춰주세요.",
    "tooSlow": "호흡이 길어졌어요. 다음 문장으로 조금 더 빨리 넘어가도 좋습니다.",
    "longSilence": "침묵이 길었습니다. 준비한 연결 문장으로 바로 회복해보세요.",
    "unclear": "문장 끝이 살짝 흐립니다. 끝 음절을 조금만 더 또렷하게 닫아주세요.",
    "offScript": "대본 핵심어가 약해졌어요. 주제어를 한 번 다시 꺼내주세요.",
}


def choose_audience_agent(session_id: str, situation: str) -> str:
    state = audience_agent_sessions.setdefault(session_id, {"last_agent": "", "last_call_at": 0.0, "llm_calls": 0})
    pool = SITUATION_AGENT_POOL.get(situation, list(AUDIENCE_AGENTS))
    if random.random() < 0.25:
        pool = list(AUDIENCE_AGENTS)
    candidates = [name for name in pool if name != state.get("last_agent")] or pool
    name = random.choice(candidates)
    state["last_agent"] = name
    return name


def requested_audience_agent(request: AudienceChatRequest) -> str | None:
    if not request.audience_name:
        return None
    candidate = request.audience_name.strip()
    return candidate if candidate in AUDIENCE_AGENTS else None


def fallback_audience_chat(session_id: str, request: AudienceChatRequest, reason: str = "fallback") -> dict[str, Any]:
    agent_name = requested_audience_agent(request) or choose_audience_agent(session_id, request.situation)
    agent = AUDIENCE_AGENTS[agent_name]
    text = (
        agent["fallback"].get(request.situation)
        or SITUATION_FALLBACK_TEXT.get(request.situation)
        or agent["fallback"].get("default", "좋아요. 핵심 문장만 더 선명하게 이어가 보세요.")
    )
    return {
        "id": f"{int(time.time() * 1000)}-{request.situation}-{agent_name}",
        "name": agent_name,
        "role": agent["role"],
        "text": text,
        "reaction": request.reaction or SITUATION_REACTIONS.get(request.situation, agent["reaction"]),
        "source": reason,
        "model": None,
    }


def should_call_audience_llm(session_id: str, force: bool = False) -> tuple[bool, str]:
    state = audience_agent_sessions.setdefault(session_id, {"last_agent": "", "last_call_at": 0.0, "llm_calls": 0})
    now = time.monotonic()
    if state.get("llm_calls", 0) >= OPENAI_AUDIENCE_MAX_SESSION_CALLS:
        return False, "session_limit"
    if not force and now - state.get("last_call_at", 0.0) < OPENAI_AUDIENCE_MIN_INTERVAL_SECONDS:
        return False, "session_throttle"

    allowed, retry_after, _remaining = openai_audience_limiter.allow()
    if not allowed:
        return False, f"global_throttle:{retry_after:.1f}"

    state["last_call_at"] = now
    state["llm_calls"] = state.get("llm_calls", 0) + 1
    return True, "allowed"


def build_audience_prompt(session: SessionState, agent_name: str, request: AudienceChatRequest) -> str:
    agent = AUDIENCE_AGENTS[agent_name]
    transcript_excerpt = compact_text(request.transcript, 700)
    current_spoken_excerpt = compact_text(request.current_excerpt, 220) or recent_excerpt(request.transcript, 180)
    script_excerpt = compact_text(session.script, 700)
    return json.dumps(
        {
            "agent": {
                "name": agent_name,
                "role": agent["role"],
                "requested_role": request.audience_role or agent["role"],
                "personality_prompt": agent["prompt"],
                "state_sensitivity": AUDIENCE_AGENT_STATE_CRITERIA.get(agent_name, ""),
            },
            "task": (
                "당신은 발표장을 보고 있는 실제 참석자입니다. 한국어 채팅 한 줄만 작성하세요. "
                "반드시 발표자를 '발표자님'이라고 부르세요. "
                "앱 분석 문장처럼 말하지 말고, 참석자가 발표자님에게 직접 말하듯 자연스럽게 말하세요. "
                "current_spoken_excerpt와 delivery_metrics를 보고 지금 들은 내용, 속도, 침묵, 대본 흐름 중 하나를 쉬운 말로 짚으세요. "
                "force_positive가 true이면 좋은 점을 먼저 말하고, 그대로 유지할 포인트를 하나만 덧붙이세요. "
                "문제가 있는 상태면 무엇이 아쉬웠는지와 바로 고칠 행동 하나를 짧게 말하세요. "
                "45자 안팎으로, 설명문·따옴표·이모지·마크다운 없이 채팅처럼 쓰세요."
            ),
            "presentation_state": {
                "situation": request.situation,
                "audience_reaction": request.reaction,
                "state_criteria": AUDIENCE_STATE_CRITERIA.get(request.situation, "현재 발표 흐름에 맞춘 청중 반응 상태입니다."),
                "elapsed_seconds": round(request.elapsed_seconds, 1),
                "words_per_minute": round(request.words_per_minute, 1),
                "syllables_per_second": round(request.syllables_per_second, 2),
                "pause_ratio_percent": round(request.pause_ratio * 100, 1),
                "silence_streak_seconds": round(request.silence_streak, 1),
                "script_overlap_percent": round(request.overlap * 100, 1),
            },
            "delivery_metrics": {
                "voice_active": request.voice_active,
                "seconds_since_recognized": round(request.seconds_since_recognized, 1),
                "words_spoken": request.words_spoken,
                "pace_hint": (
                    "빠름"
                    if request.syllables_per_second >= 7.0 or request.words_per_minute >= 180
                    else "느림"
                    if request.syllables_per_second > 0 and request.syllables_per_second < 4.5
                    else "안정"
                ),
                "pause_hint": "긴 침묵" if request.silence_streak >= 6 else "안정",
                "script_hint": "대본 핵심어 약함" if request.overlap < 0.2 and request.words_spoken >= 8 else "대본 흐름 유지",
                "force_positive": request.force_positive,
            },
            "script_excerpt": script_excerpt,
            "current_spoken_excerpt": current_spoken_excerpt,
            "spoken_excerpt": transcript_excerpt,
        },
        ensure_ascii=False,
    )


def extract_openai_text(data: dict[str, Any]) -> str:
    if data.get("output_text"):
        return str(data["output_text"]).strip()
    for output in data.get("output", []):
        for content in output.get("content", []):
            if content.get("type") in {"output_text", "text"} and content.get("text"):
                return str(content["text"]).strip()
    return ""


def mark_openai_success(kind: str) -> None:
    now = datetime.now(timezone.utc).isoformat()
    openai_runtime_state["last_live"] = True
    openai_runtime_state["last_ok_at"] = now
    openai_runtime_state["last_call_kind"] = kind
    openai_runtime_state["last_error_code"] = None
    openai_runtime_state["last_error_message"] = None


def mark_openai_error(kind: str, message: str, status_code: int | None = None) -> None:
    now = datetime.now(timezone.utc).isoformat()
    openai_runtime_state["last_live"] = False
    openai_runtime_state["last_error_at"] = now
    openai_runtime_state["last_call_kind"] = kind
    openai_runtime_state["last_error_code"] = status_code
    openai_runtime_state["last_error_message"] = compact_text(message, 500)


async def build_audience_chat(session_id: str, session: SessionState, request: AudienceChatRequest) -> dict[str, Any]:
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    if not api_key:
        mark_openai_error("audience_chat", "OPENAI_API_KEY is not configured")
        return fallback_audience_chat(session_id, request, "no_openai_key")

    can_call, reason = should_call_audience_llm(session_id, request.force)
    if not can_call:
        mark_openai_error("audience_chat", reason, 429 if "throttle" in reason or "limit" in reason else None)
        return fallback_audience_chat(session_id, request, reason)

    agent_name = requested_audience_agent(request) or choose_audience_agent(session_id, request.situation)
    agent = AUDIENCE_AGENTS[agent_name]
    payload = {
        "model": OPENAI_AUDIENCE_MODEL,
        "input": [
            {
                "role": "system",
                "content": (
                    "You create live audience feedback for a Korean presentation practice app. "
                    "Speak as the named audience member, directly to the presenter. "
                    "Always address the presenter as '발표자님' in the Korean message. "
                    "Use plain, friendly Korean that a user can immediately understand. "
                    "Output one natural Korean chat message only. No JSON, no markdown."
                ),
            },
            {"role": "user", "content": build_audience_prompt(session, agent_name, request)},
        ],
        "temperature": 0.85,
        "max_output_tokens": 80,
    }

    try:
        async with httpx.AsyncClient(timeout=8) as client:
            response = await client.post(
                "https://api.openai.com/v1/responses",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload,
            )
            response.raise_for_status()
            text = clean_user_text(extract_openai_text(response.json()))
            if not text:
                raise ValueError("empty OpenAI audience response")
            mark_openai_success("audience_chat")
            return {
                "id": f"{int(time.time() * 1000)}-{request.situation}-{agent_name}",
                "name": agent_name,
                "role": agent["role"],
                "text": compact_text(text, 90),
                "reaction": request.reaction or SITUATION_REACTIONS.get(request.situation, agent["reaction"]),
                "source": "openai",
                "model": OPENAI_AUDIENCE_MODEL,
            }
    except httpx.HTTPStatusError as exc:
        mark_openai_error("audience_chat", exc.response.text[:500] or f"HTTP {exc.response.status_code}", exc.response.status_code)
        return fallback_audience_chat(session_id, request, "openai_error")
    except Exception:
        mark_openai_error("audience_chat", "OpenAI request failed")
        return fallback_audience_chat(session_id, request, "openai_error")


def build_overlap_ratio(script_tokens: set[str], transcript: str) -> float:
    if not script_tokens:
        return 0.0
    transcript_tokens = set(tokenize(transcript))
    if not transcript_tokens:
        return 0.0
    return len(script_tokens & transcript_tokens) / len(script_tokens)


FEEDBACK_TOPIC_KEYWORDS = {
    "rhythm": ("리듬", "구간별", "변동", "일정"),
    "pace": ("속도", "말속도", "빠른", "빨라", "느린", "느려", "음절", "wpm"),
    "pause": ("침묵", "휴지", "멈춤", "쉬는", "쉼", "공백", "복구"),
    "script": ("대본", "핵심어", "키워드", "메시지", "반영"),
    "material": ("자료", "슬라이드", "시인성", "발표자료"),
    "ending": ("마무리", "결론", "감사"),
}

USER_TEXT_REPLACEMENTS = {
    "논문에서 전달력 높은 말하기로 제시된 보통 발화 속도(초당 약 6음절)에 가깝습니다.": "말 속도가 안정적이라 핵심 내용이 따라가기 좋습니다.",
    "전체 발화 중 휴지 비율이 논문에서 제시한 전달력 높은 말하기의 범위에 가깝습니다.": "쉬는 타이밍이 과하지 않아 발표 흐름이 안정적입니다.",
    "휴지 비율이 25% 이상이면 전달력이 떨어질 수 있습니다.": "쉬는 시간이 길어 흐름이 끊겨 보일 수 있습니다.",
    "논문에서 제시한 전달력 높은 말하기의 ": "",
    "논문에서 전달력 높은 말하기로 제시된 ": "",
}

UNHELPFUL_SUMMARY_MARKERS = (
    "신지영",
    "운율 중심",
    "연구를 반영",
    "criteria_basis",
    "rubric",
    "내부 기준",
)


def clean_user_text(text: Any) -> str:
    cleaned = str(text or "").strip()
    for source, replacement in USER_TEXT_REPLACEMENTS.items():
        cleaned = cleaned.replace(source, replacement)
    return re.sub(r"\s+", " ", cleaned).strip()


def feedback_topic_key(text: str) -> str:
    normalized = clean_user_text(text).lower()
    for topic, keywords in FEEDBACK_TOPIC_KEYWORDS.items():
        if any(keyword.lower() in normalized for keyword in keywords):
            return topic
    return normalized[:32]


def unique_feedback_items(items: list[Any], max_items: int | None = None) -> list[str]:
    result: list[str] = []
    seen_topics: set[str] = set()
    seen_texts: set[str] = set()
    for item in items:
        text = clean_user_text(item)
        if not text:
            continue
        exact_key = re.sub(r"\s+", "", text)
        topic_key = feedback_topic_key(text)
        if exact_key in seen_texts or topic_key in seen_topics:
            continue
        seen_texts.add(exact_key)
        seen_topics.add(topic_key)
        result.append(text)
        if max_items and len(result) >= max_items:
            break
    return result


def clean_issue_item(item: dict[str, Any]) -> dict[str, Any]:
    cleaned = dict(item)
    for key in ("title", "evidence", "spoken_excerpt", "suggestion"):
        if key in cleaned:
            cleaned[key] = clean_user_text(cleaned.get(key))
    return cleaned


def unique_issue_log(issues: list[dict[str, Any]], max_items: int = 6) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    seen_types: set[str] = set()
    for item in issues:
        issue_type = str(item.get("type") or item.get("title") or "")
        if issue_type in seen_types:
            continue
        seen_types.add(issue_type)
        result.append(clean_issue_item(item))
        if len(result) >= max_items:
            break
    return result


def default_user_summary(report: dict[str, Any]) -> str:
    score = report.get("overall_score") or 0
    if score >= 80:
        return "전체 흐름은 안정적입니다. 다음 연습에서는 강조와 마무리만 조금 더 선명하게 다듬어 보세요."
    if score >= 60:
        return "발표의 큰 흐름은 잡혀 있습니다. 속도, 쉬는 타이밍, 핵심어 전달을 조금 더 정리하면 훨씬 또렷해집니다."
    return "이번 연습에서는 흐름을 먼저 안정시키는 것이 좋습니다. 긴 침묵과 핵심어 전달을 우선 다듬어 보세요."


def user_facing_summary(report: dict[str, Any]) -> str:
    summary = clean_user_text(report.get("summary"))
    if not summary or any(marker in summary for marker in UNHELPFUL_SUMMARY_MARKERS):
        return default_user_summary(report)
    return summary


def sanitize_report_for_user(report: dict[str, Any]) -> dict[str, Any]:
    report["summary"] = user_facing_summary(report)
    report["strengths"] = unique_feedback_items(report.get("strengths") or [], 4) or ["리허설 데이터를 안정적으로 수집했습니다."]
    report["improvements"] = unique_feedback_items(report.get("improvements") or [], 6)
    report["issue_log"] = unique_issue_log(report.get("issue_log") or [], 12)
    report["timeline_log"] = [clean_issue_item(item) for item in (report.get("timeline_log") or [])[:20]]

    detailed_feedback = dict(report.get("detailed_feedback") or {})
    detailed_feedback["priority_feedback"] = unique_feedback_items(
        detailed_feedback.get("priority_feedback") or report.get("improvements") or [],
        5,
    )
    detailed_feedback["practice_plan"] = unique_feedback_items(detailed_feedback.get("practice_plan") or [], 5)
    if "coach_note" in detailed_feedback:
        detailed_feedback["coach_note"] = clean_user_text(detailed_feedback.get("coach_note"))
    report["detailed_feedback"] = detailed_feedback
    return report


def merge_report_with_fallback(
    report: dict[str, Any],
    fallback_report: dict[str, Any],
    transcript: str,
) -> dict[str, Any]:
    merged = dict(fallback_report)
    merged.update(report or {})

    merged["pace"] = fallback_report.get("pace")
    merged["silence"] = fallback_report.get("silence")
    merged["rhythm"] = fallback_report.get("rhythm")
    merged["delivery_match"] = fallback_report.get("delivery_match")
    merged["analysis_meta"] = fallback_report.get("analysis_meta")
    merged["speech_habits"] = fallback_report.get("speech_habits")
    merged["keyword_feedback"] = fallback_report.get("keyword_feedback")
    merged["presentation_material"] = fallback_report.get("presentation_material")
    merged["audience_reactions"] = fallback_report.get("audience_reactions")
    merged["criteria_basis"] = fallback_report.get("criteria_basis")
    merged["issue_log"] = report.get("issue_log") or fallback_report.get("issue_log") or []
    merged["timeline_log"] = report.get("timeline_log") or fallback_report.get("timeline_log") or []
    merged["reference_video"] = report.get("reference_video") or fallback_report.get("reference_video")
    merged["reference_comparison"] = report.get("reference_comparison") or fallback_report.get("reference_comparison")

    analysis_meta = fallback_report.get("analysis_meta") or {}
    cautious_mode = analysis_meta.get("level") != "full" or (fallback_report.get("overall_score") or 0) < 65
    if cautious_mode:
        merged["strengths"] = fallback_report.get("strengths") or []
    else:
        merged["strengths"] = unique_feedback_items(
            (fallback_report.get("strengths") or []) + (report.get("strengths") or []),
            4,
        )

    merged["improvements"] = unique_feedback_items(
        (report.get("improvements") or []) + (fallback_report.get("improvements") or []),
        6,
    )

    fallback_detail = fallback_report.get("detailed_feedback") or {}
    report_detail = report.get("detailed_feedback") or {}
    merged["detailed_feedback"] = {
        "priority_feedback": unique_feedback_items(
            (report_detail.get("priority_feedback") or []) + (fallback_detail.get("priority_feedback") or []),
            5,
        ),
        "practice_plan": unique_feedback_items(
            (report_detail.get("practice_plan") or []) + (fallback_detail.get("practice_plan") or []),
            5,
        ),
        "coach_note": report_detail.get("coach_note") or fallback_detail.get("coach_note") or "",
    }
    merged["transcript_full"] = transcript.strip()
    merged["analysis_basis"] = {
        "analysis_source": "gemini_plus_heuristic" if merged.get("used_gemini") else "heuristic",
        "timeline_count": len(merged.get("timeline_log") or []),
        "issue_count": len(merged.get("issue_log") or []),
        "speech_samples": analysis_meta.get("speech_samples", 0),
        "spoken_words": analysis_meta.get("spoken_words", 0),
        "duration_seconds": analysis_meta.get("duration_seconds", 0),
    }
    return merged


def build_issue_log(session: SessionState, transcript: str) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    previous_transcript = ""
    previous_longest_silence = 0.0
    script_tokens = set(tokenize(session.script))

    def add_issue(
        sample: MetricSample,
        issue_type: str,
        severity: str,
        title: str,
        evidence: str,
        suggestion: str,
        excerpt: str,
        metric: dict[str, Any],
    ) -> None:
        duplicate_key = (issue_type, format_seconds(sample.elapsed_seconds))
        if any((item["type"], item["time"]) == duplicate_key for item in issues):
            return
        issues.append(
            {
                "time": duplicate_key[1],
                "elapsed_seconds": round(sample.elapsed_seconds, 1),
                "type": issue_type,
                "severity": severity,
                "title": title,
                "evidence": evidence,
                "spoken_excerpt": excerpt or "해당 구간에서 새로 인식된 문장이 거의 없었습니다.",
                "metric": metric,
                "suggestion": suggestion,
            }
        )

    for sample in session.samples:
        current_transcript = sample.transcript.strip()
        delta = transcript_delta(previous_transcript, current_transcript)
        excerpt = recent_excerpt(delta or current_transcript)
        overlap = build_overlap_ratio(script_tokens, current_transcript)

        if sample.elapsed_seconds >= 8 and sample.syllables_per_second >= 7.2:
            add_issue(
                sample,
                "pace_fast",
                "high" if sample.syllables_per_second >= 8 else "medium",
                "말이 빨라 핵심어가 지나간 구간",
                f"초당 {sample.syllables_per_second:.1f}음절로 목표 범위(5.6-6.3)를 넘었습니다.",
                "핵심 단어 앞뒤로 반 박자 쉬고, 한 문장을 두 덩어리로 끊어 말해보세요.",
                excerpt,
                {"syllables_per_second": round(sample.syllables_per_second, 2), "target": "5.6-6.3"},
            )

        if sample.elapsed_seconds >= 12 and 0 < sample.syllables_per_second < 4.8:
            add_issue(
                sample,
                "pace_slow",
                "medium",
                "흐름이 느려진 구간",
                f"초당 {sample.syllables_per_second:.1f}음절로 목표 범위보다 낮았습니다.",
                "문장 사이 연결어를 미리 정해 두고 다음 의미 단위로 바로 넘어가세요.",
                excerpt,
                {"syllables_per_second": round(sample.syllables_per_second, 2), "target": "5.6-6.3"},
            )

        new_silence = max(0, sample.longest_silence_seconds - previous_longest_silence)
        if sample.longest_silence_seconds >= 5 and new_silence >= 2:
            add_issue(
                sample,
                "long_silence",
                "high" if sample.longest_silence_seconds >= 8 else "medium",
                "긴 침묵이 생긴 구간",
                f"최장 침묵이 {sample.longest_silence_seconds:.0f}초까지 늘었습니다.",
                "막히는 지점에는 '다음으로는', '핵심은' 같은 복구 문장을 준비해두세요.",
                recent_excerpt(current_transcript),
                {"longest_silence_seconds": round(sample.longest_silence_seconds, 1), "risk": "5초 이상"},
            )

        if sample.elapsed_seconds >= 20 and sample.pause_ratio >= 0.28:
            add_issue(
                sample,
                "pause_ratio",
                "high",
                "쉬는 시간이 누적된 구간",
                f"전체 시간 중 휴지 비율이 {sample.pause_ratio * 100:.1f}%로 위험 기준(25%+)을 넘었습니다.",
                "문장을 완전히 멈추기보다 의미 단위 끝에서 짧게 끊어 말하는 방식으로 바꿔보세요.",
                excerpt,
                {"pause_ratio_percent": round(sample.pause_ratio * 100, 1), "risk_ratio_percent": 25},
            )

        if sample.elapsed_seconds >= 24 and current_transcript and overlap < 0.18:
            add_issue(
                sample,
                "off_script",
                "medium",
                "대본 핵심어와 멀어진 구간",
                f"현재까지 대본 핵심어 반영도가 {overlap * 100:.0f}%에 머물렀습니다.",
                "슬라이드 제목이나 대본의 중심 키워드를 다시 말해 메시지를 회수하세요.",
                excerpt,
                {"script_overlap_percent": round(overlap * 100), "expected_min_percent": 25},
            )

        previous_transcript = current_transcript
        previous_longest_silence = max(previous_longest_silence, sample.longest_silence_seconds)

    if not issues and transcript:
        last_sample = session.samples[-1] if session.samples else None
        issues.append(
            {
                "time": format_seconds(last_sample.elapsed_seconds if last_sample else 0),
                "elapsed_seconds": round(last_sample.elapsed_seconds if last_sample else 0, 1),
                "type": "review",
                "severity": "low",
                "title": "큰 위험 구간은 적었습니다",
                "evidence": "속도와 침묵 지표에서 강한 경고가 감지되지 않았습니다.",
                "spoken_excerpt": recent_excerpt(transcript),
                "metric": {},
                "suggestion": "다음 연습에서는 강조하고 싶은 문장 2-3개를 정해 억양과 쉼을 더 의식해보세요.",
            }
        )

    severity_order = {"high": 0, "medium": 1, "low": 2}
    issues.sort(key=lambda item: (severity_order.get(item["severity"], 3), item["elapsed_seconds"]))
    return unique_issue_log(issues, 6)


def classify_timeline_window(
    script_tokens: set[str],
    samples: list[MetricSample],
    previous_transcript: str,
) -> dict[str, Any]:
    start_sample = samples[0]
    end_sample = samples[-1]
    transcript = end_sample.transcript.strip()
    overlap = build_overlap_ratio(script_tokens, transcript)
    avg_sps = mean(sample.syllables_per_second for sample in samples) if samples else 0
    avg_wpm = mean(sample.words_per_minute for sample in samples) if samples else 0
    max_pause_ratio = max(sample.pause_ratio for sample in samples) if samples else 0
    max_long_silence = max(sample.longest_silence_seconds for sample in samples) if samples else 0
    speech_detected = any(sample.speech_detected for sample in samples)
    delta = transcript_delta(previous_transcript, transcript)
    excerpt = recent_excerpt(delta or transcript)

    title = "안정적으로 전달한 구간"
    evidence = (
        f"{format_seconds(start_sample.elapsed_seconds)}-{format_seconds(end_sample.elapsed_seconds)} 동안 "
        f"초당 {avg_sps:.1f}음절, 휴지 비율 {max_pause_ratio * 100:.0f}%, 대본 반영 {overlap * 100:.0f}%입니다."
    )
    suggestion = "지금 구간의 속도와 문장 연결 방식을 다음 핵심 구간에도 유지해 보세요."
    severity = "low"
    log_type = "steady"

    if not speech_detected:
        title = "음성 인식이 거의 잡히지 않은 구간"
        evidence = (
            f"{format_seconds(start_sample.elapsed_seconds)}-{format_seconds(end_sample.elapsed_seconds)} 구간에서 "
            "음성 인식 데이터가 거의 없었습니다."
        )
        suggestion = "마이크 거리와 브라우저 음성 인식 상태를 먼저 확인해 보세요."
        severity = "medium"
        log_type = "recognition_gap"
    elif max_long_silence >= 8 or max_pause_ratio >= 0.32:
        title = "침묵이 길어 흐름이 끊긴 구간"
        evidence = (
            f"최장 침묵 {max_long_silence:.1f}초, 휴지 비율 {max_pause_ratio * 100:.0f}%로 "
            "청중 입장에서 멈춘 느낌이 크게 날 수 있습니다."
        )
        suggestion = "다음 문장으로 넘어갈 때 바로 꺼낼 연결 문장을 미리 준비해 두세요."
        severity = "high"
        log_type = "long_silence"
    elif avg_sps >= 7.1 or avg_wpm >= 175:
        title = "속도가 빨라 정보가 밀릴 수 있는 구간"
        evidence = f"평균 초당 {avg_sps:.1f}음절, 분당 {avg_wpm:.0f}단어 수준으로 목표 범위를 넘었습니다."
        suggestion = "문장 끝마다 짧게 끊고 핵심 명사에만 힘을 실어 전달해 보세요."
        severity = "medium" if avg_sps < 8 else "high"
        log_type = "pace_fast"
    elif 0 < avg_sps < 4.8:
        title = "속도가 느려 리듬이 처진 구간"
        evidence = f"평균 초당 {avg_sps:.1f}음절로 전체 리듬이 느슨하게 들릴 수 있습니다."
        suggestion = "문장 첫 단어를 더 또렷하게 시작해서 전개 속도를 조금 끌어올려 보세요."
        severity = "medium"
        log_type = "pace_slow"
    elif overlap < 0.2 and len(tokenize(transcript)) >= 10:
        title = "대본 핵심어와 멀어진 구간"
        evidence = f"이 시점 대본 반영률이 {overlap * 100:.0f}% 수준으로 내려가 핵심 메시지가 흐려질 수 있습니다."
        suggestion = "슬라이드 제목이나 결론 키워드를 다시 말해 흐름을 붙잡아 보세요."
        severity = "medium"
        log_type = "off_script"
    elif avg_sps >= 5.5 and avg_sps <= 6.4 and max_pause_ratio >= 0.1 and max_pause_ratio <= 0.22 and overlap >= 0.35:
        title = "속도와 호흡이 잘 맞은 구간"
        evidence = (
            f"평균 초당 {avg_sps:.1f}음절, 휴지 비율 {max_pause_ratio * 100:.0f}%, "
            f"대본 반영 {overlap * 100:.0f}%로 안정적이었습니다."
        )
        suggestion = "이 구간을 기준 템포로 삼고 비슷한 말하기 리듬을 반복 연습해 보세요."
        severity = "low"
        log_type = "strong_segment"

    return {
        "time": f"{format_seconds(start_sample.elapsed_seconds)}-{format_seconds(end_sample.elapsed_seconds)}",
        "elapsed_seconds": round(end_sample.elapsed_seconds, 1),
        "type": log_type,
        "severity": severity,
        "title": title,
        "evidence": evidence,
        "spoken_excerpt": excerpt or "이 구간에서 새로 인식된 발표 문장이 많지 않았습니다.",
        "metric": {
            "average_syllables_per_second": round(avg_sps, 2),
            "average_words_per_minute": round(avg_wpm, 1),
            "pause_ratio_percent": round(max_pause_ratio * 100, 1),
            "longest_silence_seconds": round(max_long_silence, 1),
            "script_overlap_percent": round(overlap * 100),
        },
        "suggestion": suggestion,
    }


def build_timeline_log(session: SessionState) -> list[dict[str, Any]]:
    if not session.samples:
        return []

    script_tokens = set(tokenize(session.script))
    logs: list[dict[str, Any]] = []
    previous_transcript = ""
    window_size = 2

    for start in range(0, len(session.samples), window_size):
        window = session.samples[start : start + window_size]
        if not window:
            continue
        logs.append(classify_timeline_window(script_tokens, window, previous_transcript))
        previous_transcript = window[-1].transcript.strip()

    return logs


def build_keyword_feedback(script: str, transcript: str) -> dict[str, Any]:
    script_tokens = [token for token in tokenize(script) if len(token) > 1]
    spoken_tokens = set(tokenize(transcript))
    seen: set[str] = set()
    important_tokens = []
    for token in script_tokens:
        if token in seen:
            continue
        seen.add(token)
        important_tokens.append(token)

    covered = [token for token in important_tokens if token in spoken_tokens]
    missed = [token for token in important_tokens if token not in spoken_tokens]
    return {
        "covered_keywords": covered[:12],
        "missed_keywords": missed[:12],
        "coverage_percent": round((len(covered) / len(important_tokens)) * 100) if important_tokens else 0,
    }


def build_detailed_feedback(report: dict[str, Any], issue_log: list[dict[str, Any]]) -> dict[str, Any]:
    pace = report["pace"]
    silence = report["silence"]
    rhythm = report["rhythm"]
    delivery = report["delivery_match"]
    analysis_meta = report.get("analysis_meta") or {}

    priorities = []
    if issue_log:
        primary_issue = issue_log[0]
        priorities.append(
            f"{primary_issue['time']} 구간에서 {primary_issue['title']}이 확인됐습니다. {primary_issue['suggestion']}"
        )
    if silence["pause_ratio_percent"] >= 25 or silence["longest_seconds"] >= 6:
        priorities.append("침묵 복구 문장을 먼저 준비하세요. 멈췄을 때 바로 이어갈 문장 하나가 전체 전달력을 크게 올립니다.")
    if pace["syllables_per_second"] > 6.8:
        priorities.append("속도를 낮추는 것이 1순위입니다. 핵심어 앞뒤에 짧은 쉼을 넣어 청중이 따라올 시간을 주세요.")
    elif pace["syllables_per_second"] < 5.0:
        priorities.append("발화 흐름을 조금 더 밀도 있게 가져가세요. 문장 사이 공백을 줄이고 다음 문장 첫 단어를 미리 떠올리면 좋습니다.")
    if delivery["similarity_percent"] < 55:
        priorities.append("대본의 핵심 키워드가 충분히 전달되지 않았습니다. 슬라이드 제목, 문제, 해결책 키워드는 반드시 소리 내어 반복하세요.")
    if rhythm["rate_variability"] > 1.2:
        priorities.append("구간별 속도 차이가 큽니다. 문장 안에서는 일정하게, 문장 끝에서만 쉬는 리듬을 연습하세요.")

    if not priorities:
        priorities.append("지표가 안정적입니다. 다음에는 강조 문장과 마무리 문장의 힘을 더 살려 완성도를 높여보세요.")

    practice_plan = [
        "문제가 표시된 구간의 발화 문장을 다시 읽고, 쉼표를 넣을 위치를 표시합니다.",
        "같은 구간만 2회 반복하며 첫 번째는 천천히, 두 번째는 실제 발표 속도로 말합니다.",
        "마지막 20초에는 결론과 감사 인사를 끊기지 않게 말하는지 확인합니다.",
    ]

    if issue_log:
        practice_plan.insert(1, f"가장 먼저 볼 구간은 {issue_log[0]['time']}의 '{issue_log[0]['title']}'입니다.")

    return {
        "priority_feedback": unique_feedback_items(priorities, 5),
        "practice_plan": practice_plan,
        "coach_note": (
            f"아래 내용은 총 {analysis_meta.get('spoken_words', 0)}개의 인식 단어와 "
            f"{analysis_meta.get('speech_samples', 0)}개 인식 구간, 그리고 속도/침묵 기록을 바탕으로 만든 근거입니다."
        ),
    }


def count_phrase_occurrences(tokens: list[str], phrases: list[str]) -> dict[str, int]:
    counts: dict[str, int] = {}
    for phrase in phrases:
        phrase_tokens = tokenize(phrase)
        if not phrase_tokens:
            continue
        if len(phrase_tokens) == 1:
            count = sum(1 for token in tokens if token == phrase_tokens[0])
        else:
            count = 0
            last_index = len(tokens) - len(phrase_tokens) + 1
            for index in range(max(0, last_index)):
                if tokens[index : index + len(phrase_tokens)] == phrase_tokens:
                    count += 1
        if count:
            counts[phrase] = count
    return counts


def build_speech_habits(transcript: str) -> dict[str, Any]:
    normalized = re.sub(r"\s+", " ", transcript).strip()
    tokens = tokenize(normalized)
    filler_counts = count_phrase_occurrences(
        tokens,
        ["어", "음", "그", "약간", "조금", "좀", "사실", "일단", "그러니까", "뭔가"],
    )
    vague_counts = count_phrase_occurrences(
        tokens,
        ["이런", "그런", "이 부분", "저 부분", "뭔가", "약간", "같은 경우"],
    )
    repair_counts = count_phrase_occurrences(
        tokens,
        ["아 아니", "다시 말하면", "정정하자면", "그러니까 다시", "아 잠시만"],
    )

    token_counter: dict[str, int] = {}
    for token in tokens:
        if len(token) <= 1:
            continue
        token_counter[token] = token_counter.get(token, 0) + 1

    repeated_tokens = [
        {"token": token, "count": count}
        for token, count in sorted(token_counter.items(), key=lambda item: (-item[1], item[0]))
        if count >= 3
    ][:5]

    sentence_endings = {
        key: value
        for key, value in {
            "습니다": normalized.count("습니다"),
            "같습니다": normalized.count("같습니다"),
            "거든요": normalized.count("거든요"),
            "해요": normalized.count("해요"),
        }.items()
        if value
    }

    filler_total = sum(filler_counts.values())
    vague_total = sum(vague_counts.values())
    repair_total = sum(repair_counts.values())
    token_total = max(1, len(tokens))

    notes = []
    if filler_total >= 4:
        notes.append("추임새가 자주 들어가 문장 전달력이 흐려질 수 있습니다.")
    if vague_total >= 3:
        notes.append("모호한 지시 표현이 반복돼 핵심 메시지가 흐려질 수 있습니다.")
    if repair_total >= 2:
        notes.append("말을 중간에 고쳐 말하는 패턴이 반복됩니다.")
    if repeated_tokens:
        notes.append(f"반복 표현이 눈에 띕니다: {', '.join(item['token'] for item in repeated_tokens[:3])}")
    if not notes:
        notes.append("눈에 띄는 말 습관 과다는 크지 않았습니다.")

    return {
        "filler_counts": filler_counts,
        "filler_total": filler_total,
        "filler_ratio_percent": round(filler_total / token_total * 100, 1),
        "vague_counts": vague_counts,
        "vague_total": vague_total,
        "repair_counts": repair_counts,
        "repair_total": repair_total,
        "repeated_tokens": repeated_tokens,
        "sentence_endings": sentence_endings,
        "notes": notes[:4],
    }


def build_analysis_meta(samples: list[MetricSample], transcript: str) -> dict[str, Any]:
    spoken_words = len(tokenize(transcript))
    duration_seconds = round(samples[-1].elapsed_seconds, 1) if samples else 0.0
    speech_samples = len([sample for sample in samples if sample.speech_detected])

    if duration_seconds < 15 or spoken_words < 20 or speech_samples < 2:
        level = "insufficient"
    elif duration_seconds < 40 or spoken_words < 60 or speech_samples < 5:
        level = "preliminary"
    else:
        level = "full"

    return {
        "level": level,
        "duration_seconds": duration_seconds,
        "spoken_words": spoken_words,
        "speech_samples": speech_samples,
        "score_visible": level != "insufficient",
        "summary_label": {
            "insufficient": "분석 보류",
            "preliminary": "예비 분석",
            "full": "정식 분석",
        }[level],
    }


def build_stt_detailed_feedback(report: dict[str, Any], issue_log: list[dict[str, Any]]) -> dict[str, Any]:
    pace = report["pace"]
    silence = report["silence"]
    rhythm = report["rhythm"]
    habits = report.get("speech_habits") or {}
    analysis_meta = report.get("analysis_meta") or {}

    if analysis_meta.get("level") == "insufficient":
        return {
            "priority_feedback": ["말한 내용이 아직 짧아서, 먼저 기본 흐름을 더 모아보는 게 좋습니다."],
            "practice_plan": [
                "최소 30초 이상 한 흐름으로 말해 보세요.",
                "핵심 문장 3개를 끊지 않고 연결해서 한 번 더 말해 보세요.",
            ],
            "coach_note": "지금 결과는 말한 내용이 아직 적어서, 가볍게 참고하는 정도로 보면 좋습니다.",
        }

    priorities = []
    if silence["pause_ratio_percent"] >= 25 or silence["longest_seconds"] >= 6:
        priorities.append("잠깐 멈추는 구간이 조금 있습니다. 다음 문장으로 넘어갈 연결 문장을 미리 준비해 두세요.")
    if pace["syllables_per_second"] > 6.8:
        priorities.append("말 속도가 조금 빠릅니다. 핵심 정보가 밀리지 않도록 문장 끝을 또렷하게 전해보세요.")
    elif pace["syllables_per_second"] < 5.0:
        priorities.append("말 속도가 다소 느려 전체 흐름이 처질 수 있습니다. 첫 문장을 좀 더 힘 있게 시작해 보세요.")
    if rhythm["rate_variability"] > 1.2:
        priorities.append("구간별 속도 차이가 커서 흐름이 흔들립니다. 문장 사이 호흡을 조금 더 일정하게 맞춰보세요.")
    if habits.get("filler_total", 0) >= 4:
        priorities.append("추임새가 자주 들립니다. 문장을 시작하기 전에 핵심어를 먼저 떠올리고 바로 말해보세요.")
    if habits.get("vague_total", 0) >= 3:
        priorities.append("표현이 조금 두루뭉술하게 들립니다. '이 부분' 대신 구체적인 대상이나 숫자를 바로 말해보세요.")

    if not priorities:
        priorities.append("속도와 호흡이 비교적 안정적입니다. 강조할 핵심 문장만 더 또렷하게 세워보세요.")

    practice_plan = [
        "타임라인 로그에서 가장 흔들린 구간 1개를 골라 같은 내용으로 2번만 다시 말해보세요.",
        "한 번은 속도만, 한 번은 추임새 제거만 신경 써서 비교 연습해보세요.",
    ]

    if issue_log:
        practice_plan.insert(1, f"먼저 손볼 구간은 {issue_log[0]['time']}의 '{issue_log[0]['title']}'입니다.")
    if habits.get("filler_total", 0) >= 4:
        practice_plan.append("추임새가 나온 문장을 그대로 적고, 같은 뜻을 더 짧고 단정하게 다시 말해보세요.")
    elif habits.get("vague_total", 0) >= 3:
        practice_plan.append("모호한 표현을 구체 명사로 바꿔 같은 문장을 다시 말해보세요.")
    else:
        practice_plan.append("비교적 안정적인 편이니, 핵심 문장 강조만 따로 연습해도 좋습니다.")

    return {
        "priority_feedback": priorities[:5],
        "practice_plan": practice_plan[:4],
        "coach_note": "이번 리포트는 말한 흐름과 구간별 기록을 바탕으로 정리한 결과입니다.",
    }


def polish_user_report_text(value: Any) -> Any:
    if isinstance(value, str):
        replacements = {
            "STT": "말한 내용",
            "transcript": "말한 내용",
            "sample": "기록",
            "analysis_meta": "분석 단계",
            "speech_habits": "말 습관",
            "spoken_words": "말한 단어 수",
            "speech_samples": "기록된 발화 수",
        }
        polished = value
        for source, target in replacements.items():
            polished = polished.replace(source, target)
        polished = polished.replace("기준으로 만든 결과입니다.", "기준으로 정리한 결과입니다.")
        polished = polished.replace("기준입니다.", "기준으로 살펴봤습니다.")
        return polished
    if isinstance(value, list):
        return [polish_user_report_text(item) for item in value]
    if isinstance(value, dict):
        return {key: polish_user_report_text(item) for key, item in value.items()}
    return value


def presentation_criteria() -> dict[str, Any]:
    return {
        "source": "신지영, 「소통과 공감을 위한 전달력 높은 말하기의 언어학적 조건(1): 운율적 측면을 중심으로」",
        "core_principle": "정보 전달형 공적 말하기에서는 말하기 속도, 화법, 쉬는 타이밍, 강조 방식이 청중의 이해와 발표자의 인상에 큰 영향을 준다.",
        "targets": {
            "speech_rate": "말하기 속도는 한국어 기준 초당 약 5.6-6.3음절을 안정적인 전달 속도 범위로 본다.",
            "speaking_style": "화법은 문장이 지나치게 길거나 딱딱하지 않고, 청중이 따라가기 쉬운 설명형 흐름인지 본다.",
            "pause_timing": "쉬는 타이밍은 중요한 의미 단위 뒤에 짧은 여백이 생기는지 본다.",
            "emphasis": "강조 방식은 핵심어를 반복하거나 속도와 억양 변화로 분명히 드러내는지 본다.",
        },
        "rubric_items": ["말하기 속도", "화법", "쉬는 타이밍", "강조 방식"],
    }

def sample_indices(total: int, limit: int) -> list[int]:
    if total <= 0 or limit <= 0:
        return []
    if total <= limit:
        return list(range(total))
    if limit == 1:
        return [0]
    step = (total - 1) / (limit - 1)
    return sorted({min(total - 1, round(step * index)) for index in range(limit)})


def compact_text(text: str, max_chars: int = 220) -> str:
    normalized = re.sub(r"\s+", " ", (text or "")).strip()
    if len(normalized) <= max_chars:
        return normalized
    return f"{normalized[: max_chars - 1].rstrip()}..."


def compact_issue_log(issue_log: list[dict[str, Any]], max_items: int = 5) -> list[dict[str, Any]]:
    compacted = []
    for item in issue_log[:max_items]:
        compacted.append(
            {
                "time": item.get("time"),
                "severity": item.get("severity"),
                "type": item.get("type"),
                "title": item.get("title"),
                "evidence": compact_text(item.get("evidence", ""), 120),
                "spoken_excerpt": compact_text(item.get("spoken_excerpt", ""), 120),
                "suggestion": compact_text(item.get("suggestion", ""), 120),
            }
        )
    return compacted


def build_sample_summary(samples: list[MetricSample]) -> dict[str, Any]:
    if not samples:
        return {
            "sample_count": 0,
            "recent_samples": [],
            "average_words_per_minute": 0,
            "average_syllables_per_second": 0,
            "max_pause_ratio_percent": 0,
            "max_longest_silence_seconds": 0,
            "reaction_counts": {},
        }

    recent_count = min(18, len(samples))
    recent_samples = []
    for index in sample_indices(len(samples), recent_count):
        sample = samples[index]
        recent_samples.append(
            {
                "elapsed_seconds": round(sample.elapsed_seconds, 1),
                "words_per_minute": round(sample.words_per_minute, 1),
                "syllables_per_second": round(sample.syllables_per_second, 2),
                "articulation_syllables_per_second": round(sample.articulation_syllables_per_second, 2),
                "pause_ratio_percent": round(sample.pause_ratio * 100, 1),
                "longest_silence_seconds": round(sample.longest_silence_seconds, 1),
                "reaction": sample.reaction,
                "speech_detected": sample.speech_detected,
                "transcript_excerpt": compact_text(sample.transcript, 120),
            }
        )

    reaction_counts: dict[str, int] = {}
    for sample in samples:
        reaction_counts[sample.reaction] = reaction_counts.get(sample.reaction, 0) + 1

    wpms = [sample.words_per_minute for sample in samples if sample.words_per_minute > 0]
    sps_values = [sample.syllables_per_second for sample in samples if sample.syllables_per_second > 0]

    return {
        "sample_count": len(samples),
        "duration_seconds": round(samples[-1].elapsed_seconds, 1),
        "recent_samples": recent_samples,
        "average_words_per_minute": round(mean(wpms), 1) if wpms else 0,
        "average_syllables_per_second": round(mean(sps_values), 2) if sps_values else 0,
        "max_pause_ratio_percent": round(max(sample.pause_ratio for sample in samples) * 100, 1),
        "max_longest_silence_seconds": round(max(sample.longest_silence_seconds for sample in samples), 1),
        "reaction_counts": reaction_counts,
    }


def build_ai_prompt_payload(session: SessionState, fallback_report: dict[str, Any]) -> dict[str, Any]:
    detailed_feedback = fallback_report.get("detailed_feedback") or {}
    return {
        "instruction": (
            "You are a Korean presentation coach. Return strict JSON only. "
            "Use the heuristic report as the source of truth and refine wording. "
            "Do not mention research sources, papers, rubrics, criteria_basis, or internal criteria in user-facing text. "
            "Avoid repeating the same evaluation topic; keep the clearest evidence when two comments cover the same issue."
        ),
        "script_excerpt": compact_text(session.script, 3500),
        "reference_video": session.reference_video,
        "sample_summary": build_sample_summary(session.samples),
        "heuristic_report": {
            "overall_score": fallback_report.get("overall_score"),
            "summary": compact_text(fallback_report.get("summary", ""), 500),
            "strengths": (fallback_report.get("strengths") or [])[:4],
            "improvements": (fallback_report.get("improvements") or [])[:5],
            "pace": fallback_report.get("pace"),
            "silence": fallback_report.get("silence"),
            "rhythm": fallback_report.get("rhythm"),
            "delivery_match": fallback_report.get("delivery_match"),
            "analysis_meta": fallback_report.get("analysis_meta"),
            "speech_habits": fallback_report.get("speech_habits"),
            "issue_log": compact_issue_log(fallback_report.get("issue_log") or []),
            "timeline_log": compact_issue_log(fallback_report.get("timeline_log") or [], 8),
            "detailed_feedback": {
                "priority_feedback": (detailed_feedback.get("priority_feedback") or [])[:4],
                "practice_plan": (detailed_feedback.get("practice_plan") or [])[:3],
                "coach_note": compact_text(detailed_feedback.get("coach_note", ""), 160),
            },
            "criteria_basis": fallback_report.get("criteria_basis"),
            "reference_video": fallback_report.get("reference_video"),
            "reference_comparison": fallback_report.get("reference_comparison"),
            "reference_speaker_comparison": fallback_report.get("reference_speaker_comparison"),
            "audience_reactions": fallback_report.get("audience_reactions", {}),
        },
        "required_schema": {
            "overall_score": "number 0-100",
            "summary": "Korean paragraph with concrete evaluation",
            "strengths": ["Korean bullet"],
            "improvements": ["Korean bullet"],
            "pace": fallback_report["pace"],
            "silence": fallback_report["silence"],
            "rhythm": fallback_report["rhythm"],
            "delivery_match": fallback_report["delivery_match"],
            "analysis_meta": fallback_report["analysis_meta"],
            "speech_habits": fallback_report["speech_habits"],
            "issue_log": fallback_report["issue_log"],
            "timeline_log": fallback_report.get("timeline_log", []),
            "detailed_feedback": fallback_report["detailed_feedback"],
            "criteria_basis": fallback_report["criteria_basis"],
            "reference_video": fallback_report.get("reference_video"),
            "reference_comparison": fallback_report.get("reference_comparison"),
            "reference_speaker_comparison": fallback_report.get("reference_speaker_comparison"),
            "audience_reactions": fallback_report["audience_reactions"],
        },
    }


def render_pdf_images_for_vision(data: bytes, page_limit: int) -> tuple[list[dict[str, Any]], list[str]]:
    if fitz is None:
        return [], ["PDF 시각 렌더링 라이브러리가 없어 텍스트 기반 분석만 가능합니다."]

    notes: list[str] = []
    images: list[dict[str, Any]] = []
    try:
        document = fitz.open(stream=data, filetype="pdf")
        for page_index in sample_indices(document.page_count, page_limit):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.7, 1.7), alpha=False)
            images.append(
                {
                    "page_index": page_index,
                    "mime_type": "image/png",
                    "data": base64.b64encode(pixmap.tobytes("png")).decode("ascii"),
                }
            )
        return images, notes
    except Exception as exc:
        return [], [f"PDF 페이지를 이미지로 렌더링하지 못했습니다: {exc}"]


def convert_pptx_to_pdf_bytes(data: bytes) -> tuple[bytes | None, list[str]]:
    notes: list[str] = []
    with tempfile.TemporaryDirectory(prefix="pptx-render-") as tmp_dir_name:
        tmp_dir = Path(tmp_dir_name)
        source_path = tmp_dir / "material.pptx"
        output_dir = tmp_dir / "out"
        output_dir.mkdir(exist_ok=True)
        source_path.write_bytes(data)
        command = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(source_path),
        ]
        try:
            subprocess.run(
                command,
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=45,
            )
        except FileNotFoundError:
            return None, ["LibreOffice가 없어 PPTX 시각 분석을 수행할 수 없습니다."]
        except subprocess.TimeoutExpired:
            return None, ["PPTX를 PDF로 변환하는 데 시간이 너무 오래 걸렸습니다."]
        except subprocess.CalledProcessError as exc:
            stderr = exc.stderr.decode("utf-8", errors="ignore").strip()
            return None, [f"PPTX를 PDF로 변환하지 못했습니다: {stderr or exc}"]

        pdf_path = output_dir / "material.pdf"
        if not pdf_path.exists():
            return None, ["PPTX를 PDF로 변환했지만 결과 파일을 찾지 못했습니다."]
        return pdf_path.read_bytes(), notes


def render_material_images_for_vision(name: str, content_type: str | None, data: bytes) -> tuple[list[dict[str, Any]], list[str]]:
    kind = file_kind(name, content_type)
    if kind == "pdf":
        return render_pdf_images_for_vision(data, MAX_VISION_PAGES_PER_FILE)
    if kind == "pptx":
        pdf_bytes, notes = convert_pptx_to_pdf_bytes(data)
        if not pdf_bytes:
            return [], notes
        rendered_images, render_notes = render_pdf_images_for_vision(pdf_bytes, MAX_VISION_PAGES_PER_FILE)
        return rendered_images, notes + render_notes
    return [], ["PDF 또는 PPTX만 시각 분석 대상으로 지원합니다."]


async def analyze_preflight_with_gemini(
    script: str,
    local_materials: list[dict[str, Any]],
    uploads: list[dict[str, Any]],
) -> dict[str, Any] | None:
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        return None

    vision_parts: list[dict[str, Any]] = []
    material_context: list[dict[str, Any]] = []
    image_budget = MAX_VISION_MATERIAL_IMAGES

    for upload, local_material in zip(uploads, local_materials):
        rendered_images, render_notes = render_material_images_for_vision(
            upload["filename"],
            upload.get("content_type"),
            upload["data"],
        )
        local_material.setdefault("notes", [])
        local_material["notes"] = list(local_material["notes"]) + render_notes

        selected_images = rendered_images[:image_budget]
        image_budget -= len(selected_images)
        material_context.append(
            {
                "filename": upload["filename"],
                "kind": local_material.get("kind"),
                "page_count": local_material.get("page_count", 0),
                "slide_count": local_material.get("slide_count", 0),
                "word_count": local_material.get("word_count", 0),
                "extracted_character_count": local_material.get("extracted_character_count", 0),
                "preview": local_material.get("preview", ""),
                "rendered_image_count": len(selected_images),
            }
        )
        for image in selected_images:
            unit_label = "page" if local_material.get("kind") == "pdf" else "slide"
            vision_parts.append(
                {"text": f"{upload['filename']} {unit_label} {image['page_index'] + 1} visual"}
            )
            vision_parts.append(
                {"inline_data": {"mime_type": image["mime_type"], "data": image["data"]}}
            )
        if image_budget <= 0:
            break

    prompt = {
        "instruction": (
            "You are a Korean presentation coach. Analyze the speech script first, then analyze the uploaded "
            "presentation materials together with that script. Return strict JSON only. Base visual judgments on the "
            "rendered slide/page images when they exist, and use extracted text only as supporting evidence."
        ),
        "script": script,
        "materials": material_context,
        "required_schema": {
            "script_feedback": {
                "score": "integer 0-100",
                "word_count": "number",
                "average_sentence_words": "number",
                "summary": "Korean short paragraph",
                "suggestions": ["Korean short sentence"],
            },
            "presentation_material": {
                "summary": "Korean short paragraph",
                "estimated_minutes": "number",
                "clarity_score": "integer 0-100",
                "consistency_score": "integer 0-100",
                "topic_fit_score": "integer 0-100",
                "overall_score": "integer 0-100",
                "notes": ["Korean short sentence"],
                "files": [
                    {
                        "filename": "string",
                        "estimated_minutes": "number",
                        "clarity_score": "integer 0-100",
                        "consistency_score": "integer 0-100",
                        "topic_fit_score": "integer 0-100",
                        "overall_score": "integer 0-100",
                        "summary": "Korean short sentence",
                        "notes": ["Korean short sentence"],
                    }
                ],
            },
        },
    }
    payload = {
        "contents": [
            {
                "parts": [{"text": json.dumps(prompt, ensure_ascii=False)}, *vision_parts],
            }
        ],
        "generationConfig": {"responseMimeType": "application/json"},
    }

    try:
        data, _remaining = await call_gemini_api(
            kind="preflight_materials",
            payload=payload,
            timeout_seconds=45,
            count_against_limit=True,
            retry_on_unavailable=1,
        )
        return json.loads(data["candidates"][0]["content"]["parts"][0]["text"])
    except Exception:
        return None


def merge_gemini_script_feedback(
    local_feedback: dict[str, Any],
    gemini_feedback: dict[str, Any] | None,
) -> dict[str, Any]:
    merged = dict(local_feedback)
    if not gemini_feedback:
        return merged

    merged.update(
        {
            "score": clamp(gemini_feedback.get("score", merged.get("score", 0)), 0, 100),
            "word_count": int(gemini_feedback.get("word_count", merged.get("word_count", 0)) or 0),
            "average_sentence_words": round(
                float(gemini_feedback.get("average_sentence_words", merged.get("average_sentence_words", 0)) or 0),
                1,
            ),
            "summary": gemini_feedback.get("summary", merged.get("summary", "")),
            "suggestions": list(gemini_feedback.get("suggestions") or merged.get("suggestions") or [])[:4],
            "analysis_source": "gemini_preflight",
        }
    )
    return merged


def merge_vision_material_feedback(
    local_materials: list[dict[str, Any]],
    vision_feedback: dict[str, Any] | None,
) -> list[dict[str, Any]]:
    if not vision_feedback:
        return local_materials

    by_name = {item.get("filename"): item for item in (vision_feedback.get("files") or [])}
    merged_materials: list[dict[str, Any]] = []
    for item in local_materials:
        merged = dict(item)
        vision_item = by_name.get(item.get("filename"))
        if vision_item:
            merged.update(
                {
                    "estimated_minutes": vision_item.get("estimated_minutes", item.get("estimated_minutes")),
                    "clarity_score": vision_item.get("clarity_score", item.get("clarity_score")),
                    "consistency_score": vision_item.get("consistency_score", item.get("consistency_score")),
                    "topic_fit_score": vision_item.get("topic_fit_score", item.get("topic_fit_score")),
                    "overall_score": vision_item.get("overall_score", item.get("overall_score")),
                    "summary": vision_item.get("summary", item.get("summary")),
                    "notes": list(item.get("notes", [])) + list(vision_item.get("notes", [])),
                    "analysis_source": "gemini_vision",
                }
            )
        merged_materials.append(merged)
    return merged_materials


def merge_material_feedback_summary(
    local_feedback: dict[str, Any],
    gemini_feedback: dict[str, Any] | None,
) -> dict[str, Any]:
    if not gemini_feedback or not local_feedback.get("uploaded"):
        return local_feedback

    merged = dict(local_feedback)
    merged.update(
        {
            "estimated_minutes": float(gemini_feedback.get("estimated_minutes", merged.get("estimated_minutes", 0)) or 0),
            "clarity_score": round(clamp(gemini_feedback.get("clarity_score", merged.get("clarity_score", 0)), 0, 100)),
            "consistency_score": round(
                clamp(gemini_feedback.get("consistency_score", merged.get("consistency_score", 0)), 0, 100)
            ),
            "topic_fit_score": round(
                clamp(gemini_feedback.get("topic_fit_score", merged.get("topic_fit_score", 0)), 0, 100)
            ),
            "overall_score": round(clamp(gemini_feedback.get("overall_score", merged.get("overall_score", 0)), 0, 100)),
            "summary": gemini_feedback.get("summary", merged.get("summary", "")),
            "notes": list(gemini_feedback.get("notes") or merged.get("notes") or [])[:8],
            "analysis_source": "gemini_preflight",
        }
    )
    return merged


def normalize_whitespace(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def split_text_units(text: str) -> list[str]:
    units = [part.strip() for part in re.split(r"\n{2,}|(?<=[.!?。！？])\s+", text) if part.strip()]
    if units:
        return units
    compact = normalize_whitespace(text)
    return [compact] if compact else []


def file_kind(filename: str, content_type: str | None) -> str:
    extension = os.path.splitext(filename.lower())[1]
    if extension == ".pdf" or (content_type and "pdf" in content_type):
        return "pdf"
    if extension == ".pptx" or (content_type and "presentation" in content_type):
        return "pptx"
    return "unknown"


def extract_pdf_text(data: bytes) -> tuple[str, int, list[str]]:
    if PdfReader is None:
        return "", 0, ["PDF 텍스트 추출 라이브러리가 설치되지 않았습니다."]

    notes: list[str] = []
    try:
        reader = PdfReader(BytesIO(data))
        texts: list[str] = []
        for index, page in enumerate(reader.pages):
            if index >= 80:
                notes.append("앞쪽 80페이지만 분석했습니다.")
                break
            try:
                texts.append(page.extract_text() or "")
            except Exception:
                texts.append("")
        return "\n".join(texts), len(reader.pages), notes
    except Exception as exc:
        return "", 0, [f"PDF를 읽는 중 문제가 발생했습니다: {exc}"]


def extract_pptx_text(data: bytes) -> tuple[str, int, list[str]]:
    if Presentation is None:
        return "", 0, ["PPTX 텍스트 추출 라이브러리가 설치되지 않았습니다."]

    notes: list[str] = []
    try:
        presentation = Presentation(BytesIO(data))
        texts: list[str] = []
        for index, slide in enumerate(presentation.slides):
            if index >= 80:
                notes.append("앞쪽 80개 슬라이드까지만 분석했습니다.")
                break
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    extracted = normalize_whitespace(shape.text_frame.text or "")
                    if extracted:
                        slide_texts.append(extracted)
            texts.append(" ".join(slide_texts))
        return "\n".join(texts), len(presentation.slides), notes
    except Exception as exc:
        return "", 0, [f"PPTX를 읽는 중 문제가 발생했습니다: {exc}"]


def weighted_keyword_overlap(script_text: str, document_text: str, lead_units: list[str]) -> float:
    script_tokens = [token for token in tokenize(script_text) if len(token) > 1]
    document_tokens = [token for token in tokenize(document_text) if len(token) > 1]
    if not script_tokens or not document_tokens:
        return 0.0

    script_set = set(script_tokens)
    lead_tokens = set()
    for unit in lead_units:
        lead_tokens.update(token for token in tokenize(unit) if len(token) > 1)

    title_hits = len(script_set & lead_tokens)
    document_hits = len(script_set & set(document_tokens))
    weighted_hits = (title_hits * 2.0) + document_hits
    weighted_total = (len(script_set) * 1.25) + max(1, len(set(document_tokens)))
    return min(1.0, weighted_hits / weighted_total * 3.2)


def analyze_material_file(script: str, name: str, content_type: str | None, data: bytes) -> dict[str, Any]:
    kind = file_kind(name, content_type)
    if len(data) > MAX_PRESENTATION_UPLOAD_BYTES:
        return {
            "filename": name,
            "content_type": content_type or "application/octet-stream",
            "kind": kind,
            "error": "파일 크기가 허용 한도를 초과했습니다.",
        }

    if kind == "pdf":
        extracted_text, page_count, notes = extract_pdf_text(data)
        unit_count = page_count or len(split_text_units(extracted_text))
    elif kind == "pptx":
        extracted_text, page_count, notes = extract_pptx_text(data)
        unit_count = page_count or len(split_text_units(extracted_text))
    else:
        return {
            "filename": name,
            "content_type": content_type or "application/octet-stream",
            "kind": kind,
            "error": "PDF 또는 PPTX만 업로드할 수 있습니다.",
        }

    extracted_text = normalize_whitespace(extracted_text)
    units = split_text_units(extracted_text)
    lead_units = units[:2]
    script_tokens = [token for token in tokenize(script) if len(token) > 1]
    material_tokens = [token for token in tokenize(extracted_text) if len(token) > 1]
    material_sets = [set(tokenize(unit)) for unit in units if tokenize(unit)]
    unit_total = max(1, unit_count or len(material_sets) or 1)
    total_word_count = len(material_tokens)
    avg_words_per_unit = total_word_count / unit_total
    longest_unit = max((len(tokens) for tokens in material_sets), default=0)
    overlap = len(set(script_tokens) & set(material_tokens)) / len(set(script_tokens)) if script_tokens and material_tokens else 0

    adjacent_overlaps: list[float] = []
    for previous, current in zip(material_sets, material_sets[1:]):
        union = previous | current
        if not union:
            continue
        adjacent_overlaps.append(len(previous & current) / len(union))

    consistency_score = round(
        clamp(
            58
            + (sum(adjacent_overlaps) / len(adjacent_overlaps) * 90 if adjacent_overlaps else 8)
            - max(0, avg_words_per_unit - 30) * 1.3
            - max(0, longest_unit - 70) * 0.4,
            35,
            100,
        )
    )
    clarity_score = round(
        clamp(
            92
            - max(0, avg_words_per_unit - 24) * 1.9
            - max(0, longest_unit - 45) * 0.65
            - (0 if extracted_text else 14),
            35,
            100,
        )
    )
    topic_alignment = weighted_keyword_overlap(script, extracted_text, lead_units)
    topic_fit_score = round(clamp((overlap * 58) + (topic_alignment * 42), 30, 100))

    estimated_minutes = round(
        max(
            max(1.0, len(script_tokens) / 130.0),
            max(1.0, unit_total * 0.45 + total_word_count / 180.0),
        )
        * 1.15,
        1,
    )

    notes = notes[:]
    if not extracted_text:
        notes.append("텍스트를 충분히 읽지 못해 이미지/도형 중심의 시인성 평가는 제한됩니다.")
    if avg_words_per_unit > 40:
        notes.append("장당 텍스트가 많아 슬라이드가 다소 붐빌 수 있습니다.")
    if overlap < 0.2:
        notes.append("대본과 자료의 공통 키워드가 적어서 주제 적합도가 낮습니다.")
    if topic_alignment >= 0.55:
        notes.append("대본 앞부분과 슬라이드 도입부의 주제가 잘 맞습니다.")
    elif topic_alignment < 0.25:
        notes.append("슬라이드 도입부와 대본의 주제 연결을 조금 더 분명하게 맞추면 좋습니다.")
    if consistency_score >= 80:
        notes.append("슬라이드 간 흐름이 비교적 안정적입니다.")

    overall_score = round((clarity_score * 0.34) + (consistency_score * 0.33) + (topic_fit_score * 0.33))
    summary = (
        f"{name} 기준 예상 발표 시간은 약 {estimated_minutes}분이고, "
        f"시인성 {clarity_score}점, 통일성 {consistency_score}점, 주제 적합도 {topic_fit_score}점으로 평가됩니다."
    )

    return {
        "filename": name,
        "content_type": content_type or "application/octet-stream",
        "kind": kind,
        "page_count": page_count if kind == "pdf" else 0,
        "slide_count": page_count if kind == "pptx" else 0,
        "extracted_character_count": len(extracted_text),
        "word_count": total_word_count,
        "unit_count": unit_total,
        "estimated_minutes": estimated_minutes,
        "clarity_score": clarity_score,
        "consistency_score": consistency_score,
        "topic_fit_score": topic_fit_score,
        "overall_score": overall_score,
        "preview": extracted_text[:280],
        "notes": notes,
        "summary": summary,
        "error": None,
    }


def build_material_feedback(script: str, materials: list[dict[str, Any]]) -> dict[str, Any]:
    if not materials:
        return {
            "uploaded": False,
            "files": [],
            "estimated_minutes": max(1.0, round(len(tokenize(script)) / 130.0, 1)),
            "clarity_score": 0,
            "consistency_score": 0,
            "topic_fit_score": 0,
            "overall_score": 0,
            "summary": "발표 자료가 업로드되지 않았습니다.",
            "notes": [],
        }

    valid_files = [item for item in materials if not item.get("error")]
    estimated_minutes = max((item.get("estimated_minutes", 0) for item in valid_files), default=0.0)
    if not estimated_minutes:
        estimated_minutes = max(1.0, round(len(tokenize(script)) / 130.0, 1))

    if valid_files:
        clarity_score = round(sum(item["clarity_score"] for item in valid_files) / len(valid_files))
        consistency_score = round(sum(item["consistency_score"] for item in valid_files) / len(valid_files))
        topic_fit_score = round(sum(item["topic_fit_score"] for item in valid_files) / len(valid_files))
        overall_score = round(sum(item["overall_score"] for item in valid_files) / len(valid_files))
    else:
        clarity_score = consistency_score = topic_fit_score = overall_score = 0

    summary = " / ".join(item["summary"] for item in valid_files[:3]) if valid_files else "업로드된 자료를 해석하지 못했습니다."
    notes = [note for item in valid_files for note in item.get("notes", [])][:8]

    return {
        "uploaded": True,
        "files": materials,
        "estimated_minutes": estimated_minutes,
        "clarity_score": clarity_score,
        "consistency_score": consistency_score,
        "topic_fit_score": topic_fit_score,
        "overall_score": overall_score,
        "summary": summary,
        "notes": notes,
    }


def build_heuristic_report(session: SessionState, transcript: str) -> dict[str, Any]:
    samples = session.samples
    last_sample = samples[-1] if samples else None
    script_tokens = set(tokenize(session.script))
    transcript_tokens = tokenize(transcript)
    transcript_set = set(transcript_tokens)
    overlap = len(script_tokens & transcript_set) / len(script_tokens) if script_tokens else 0
    material_feedback = build_material_feedback(session.script, session.materials)
    keyword_feedback = build_keyword_feedback(session.script, transcript)
    analysis_meta = build_analysis_meta(samples, transcript)
    speech_habits = build_speech_habits(transcript)

    elapsed = last_sample.elapsed_seconds if last_sample else 0
    total_silence = last_sample.silence_seconds if last_sample else 0
    longest_silence = max([sample.longest_silence_seconds for sample in samples], default=0)
    pause_ratio = last_sample.pause_ratio if last_sample and last_sample.pause_ratio else (
        total_silence / elapsed if elapsed else 0
    )

    syllables = last_sample.syllables_spoken if last_sample and last_sample.syllables_spoken else count_syllables(transcript)
    speech_rate = last_sample.syllables_per_second if last_sample and last_sample.syllables_per_second else (
        syllables / elapsed if elapsed else 0
    )
    articulation_seconds = max(1, elapsed - total_silence)
    articulation_rate = (
        last_sample.articulation_syllables_per_second
        if last_sample and last_sample.articulation_syllables_per_second
        else syllables / articulation_seconds
    )
    rate_gap = max(0, articulation_rate - speech_rate)

    wpms = [sample.words_per_minute for sample in samples if sample.words_per_minute > 0]
    avg_wpm = mean(wpms) if wpms else 0
    sps_samples = [sample.syllables_per_second for sample in samples if sample.syllables_per_second > 0]
    rate_variability = pstdev(sps_samples) if len(sps_samples) > 1 else 0
    slow_sample_ratio = (
        len([value for value in sps_samples if value < 5.0]) / len(sps_samples)
        if sps_samples
        else 0
    )

    reaction_counts: dict[str, int] = {}
    for sample in samples:
        reaction_counts[sample.reaction] = reaction_counts.get(sample.reaction, 0) + 1

    pace_score = score_distance(speech_rate, 6.0, 1.2)
    pause_score = score_distance(pause_ratio, 0.15, 0.12)
    if pause_ratio >= 0.25:
        pause_score = min(pause_score, 55)
    gap_score = score_distance(rate_gap, 1.1, 0.9)
    rhythm_score = clamp(100 - rate_variability * 22 - slow_sample_ratio * 35, 35, 100)
    habit_penalty = min(18, speech_habits["filler_total"] * 1.6 + speech_habits["vague_total"] * 1.2 + speech_habits["repair_total"] * 2.0)
    habit_score = clamp(100 - habit_penalty, 45, 100)
    overall = round(
        clamp(
            pace_score * 0.3
            + pause_score * 0.28
            + gap_score * 0.12
            + rhythm_score * 0.15
            + habit_score * 0.15,
            0,
            100,
        )
    )
    if analysis_meta["level"] == "insufficient":
        overall = 0
    elif analysis_meta["level"] == "preliminary":
        overall = round(min(overall, 79))

    strengths = []
    improvements = []
    if analysis_meta["level"] == "insufficient":
        strengths.append("짧게라도 말한 내용이 모여 기본 흐름은 확인할 수 있었습니다.")
        improvements.append("발화 길이가 짧아 신뢰도 있는 분석이 어렵습니다. 30초 이상 이어서 말해보세요.")
    if 5.6 <= speech_rate <= 6.3:
        strengths.append("말 속도가 안정적이라 핵심 내용이 따라가기 좋습니다.")
    elif speech_rate < 5.0:
        improvements.append("인식된 발화 속도가 느린 편입니다. 긴 침묵을 줄이고 다음 의미 단위로 자연스럽게 이어가 보세요.")
    else:
        improvements.append("발화 속도가 빠른 편입니다. 핵심어 앞뒤에서 짧게 쉬어 청중이 정보를 처리할 시간을 주세요.")

    if 0.10 <= pause_ratio <= 0.20:
        strengths.append("쉬는 타이밍이 과하지 않아 발표 흐름이 안정적입니다.")
    elif pause_ratio >= 0.25:
        improvements.append("휴지 비율이 25% 이상이면 전달력이 떨어질 수 있습니다. 멈춤을 줄이고 의미 단위별로 끊어 말해보세요.")
    else:
        improvements.append("휴지가 너무 적으면 정보가 밀려 들릴 수 있습니다. 중요한 문장 뒤에는 짧은 쉼을 넣어보세요.")

    if rate_gap <= 1.3:
        strengths.append("말 속도와 조음 속도의 차이가 크지 않아 흐름이 비교적 안정적입니다.")
    else:
        improvements.append("말하는 구간 자체는 괜찮지만 침묵 때문에 전체 속도가 느려집니다. 준비된 연결 문장을 활용해보세요.")

    if rate_variability > 1.2 or slow_sample_ratio > 0.35:
        improvements.append("구간별 속도 변동이 큽니다. 한 문장 안에서는 일정한 리듬을 유지하고 문장 끝에서만 쉬어보세요.")

    if speech_habits["filler_total"] >= 4:
        improvements.append("추임새가 반복됩니다. 문장을 시작하기 전에 핵심어를 먼저 정리해 보세요.")
    if speech_habits["vague_total"] >= 3:
        improvements.append("모호한 표현이 반복돼 메시지가 흐려질 수 있습니다. 구체 명사를 바로 말해보세요.")
    if speech_habits["repair_total"] == 0 and analysis_meta["level"] != "insufficient":
        strengths.append("말을 중간에 고쳐 말하는 패턴이 많지 않아 문장 흐름이 비교적 안정적입니다.")
    issue_log = build_issue_log(session, transcript)
    timeline_log = build_timeline_log(session)

    report = {
        "overall_score": overall,
        "pace": {
            "average_wpm": round(avg_wpm, 1),
            "syllables_per_second": round(speech_rate, 2),
            "articulation_syllables_per_second": round(articulation_rate, 2),
            "speech_articulation_gap": round(rate_gap, 2),
            "target_range": "5.6-6.3 syllables/sec",
            "score": round(clamp(pace_score, 0, 100)),
        },
        "silence": {
            "total_seconds": round(total_silence, 1),
            "longest_seconds": round(longest_silence, 1),
            "pause_ratio_percent": round(pause_ratio * 100, 1),
            "target_ratio_percent": "around 15",
            "risk_ratio_percent": "25+",
            "score": round(clamp(pause_score, 0, 100)),
        },
        "rhythm": {
            "rate_variability": round(rate_variability, 2),
            "slow_sample_ratio_percent": round(slow_sample_ratio * 100, 1),
            "score": round(rhythm_score),
        },
        "delivery_match": {
            "spoken_words": len(transcript_tokens),
            "spoken_syllables": syllables,
        },
        "analysis_meta": analysis_meta,
        "speech_habits": speech_habits,
        "criteria_basis": presentation_criteria(),
        "keyword_feedback": keyword_feedback,
        "presentation_material": material_feedback,
        "reference_video": session.reference_video,
        "audience_reactions": reaction_counts,
        "strengths": unique_feedback_items(strengths, 4) or ["리허설 데이터를 안정적으로 수집했습니다."],
        "improvements": unique_feedback_items(improvements, 6),
        "issue_log": issue_log,
        "timeline_log": timeline_log,
        "transcript_full": transcript.strip(),
        "analysis_basis": {
            "analysis_source": "heuristic",
            "timeline_count": len(timeline_log),
            "issue_count": len(issue_log),
            "speech_samples": analysis_meta["speech_samples"],
            "spoken_words": analysis_meta["spoken_words"],
            "duration_seconds": analysis_meta["duration_seconds"],
        },
        "summary": "이번 연습에서 바로 고칠 부분을 중심으로 리포트를 정리했습니다.",
        "used_gemini": False,
    }
    if material_feedback.get("uploaded"):
        report["summary"] = f"{report['summary']} 발표 자료의 시간과 시인성도 함께 확인했습니다."
    if session.reference_video:
        report["reference_comparison"] = build_reference_comparison(report, session.reference_video)
    if analysis_meta["level"] == "insufficient":
        report["summary"] = "말한 내용이 아직 짧아서 점수형 총평보다는 간단한 참고용 분석만 제공했습니다. 30초 이상 이어서 말하면 더 정확해집니다."
    elif analysis_meta["level"] == "preliminary":
        report["summary"] = f"{report['summary']} 다만 이번 결과는 가볍게 흐름을 점검하는 예비 결과로 보면 좋습니다."
    else:
        report["summary"] = (
            f"총 {analysis_meta['spoken_words']}개의 인식 단어와 {analysis_meta['speech_samples']}개의 인식 구간을 바탕으로 "
            f"속도, 침묵, 리듬, 말 습관을 함께 살펴봤습니다."
        )
        if issue_log:
            report["summary"] += f" 가장 먼저 눈에 띈 부분은 {issue_log[0]['title']}입니다."
    report["detailed_feedback"] = build_stt_detailed_feedback(report, issue_log)

    if session.reference_video:
        report["reference_video"] = session.reference_video
        report["reference_comparison"] = build_reference_comparison(report, session.reference_video)

    user_reference_profile = build_user_profile_from_samples(samples)
    if user_reference_profile:
        report["reference_speaker_comparison"] = compare_with_reference(user_reference_profile)

    return polish_user_report_text(sanitize_report_for_user(report))


async def ask_gemini_for_report(session: SessionState, fallback_report: dict[str, Any]) -> dict[str, Any]:
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        return polish_user_report_text(sanitize_report_for_user(fallback_report))

    allowed, retry_after, remaining = gemini_limiter.allow()
    if not allowed:
        fallback_report["summary"] = "AI 분석이 잠시 막혀 기본 리포트로 정리했습니다. 지금 바로 고칠 부분부터 확인해 주세요."
        fallback_report["used_gemini"] = False
        fallback_report["gemini_limits"] = {
            "rate_limited": True,
            "retry_after_seconds": round(retry_after, 1),
            "remaining_calls": remaining,
        }
        return polish_user_report_text(sanitize_report_for_user(fallback_report))

    model = get_gemini_model()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    prompt = {
        "instruction": (
            "You are a Korean presentation coach. Return strict JSON only. "
            "Do not mention research sources, papers, rubrics, criteria_basis, or internal criteria in user-facing text. "
            "Avoid repeating the same evaluation topic; keep the clearest evidence when two comments cover the same issue."
        ),
        "script": session.script,
        "reference_video": session.reference_video,
        "samples": [sample.model_dump() for sample in session.samples[-80:]],
        "heuristic_report": fallback_report,
        "required_schema": {
            "overall_score": "number 0-100",
            "summary": "Korean paragraph with concrete evaluation",
            "strengths": ["Korean bullet"],
            "improvements": ["Korean bullet"],
            "pace": fallback_report["pace"],
            "silence": fallback_report["silence"],
            "rhythm": fallback_report["rhythm"],
            "script": fallback_report["script"],
            "delivery_match": fallback_report["delivery_match"],
            "keyword_feedback": fallback_report["keyword_feedback"],
            "issue_log": fallback_report["issue_log"],
            "detailed_feedback": fallback_report["detailed_feedback"],
            "criteria_basis": fallback_report["criteria_basis"],
            "presentation_material": fallback_report.get("presentation_material", {}),
            "reference_video": fallback_report.get("reference_video"),
            "reference_comparison": fallback_report.get("reference_comparison"),
            "reference_speaker_comparison": fallback_report.get("reference_speaker_comparison"),
            "audience_reactions": fallback_report["audience_reactions"],
            "reference_video": fallback_report.get("reference_video"),
            "reference_comparison": fallback_report.get("reference_comparison"),
        },
    }

    payload = {
        "contents": [{"parts": [{"text": json.dumps(prompt, ensure_ascii=False)}]}],
        "generationConfig": {"responseMimeType": "application/json"},
    }

    try:
        async with httpx.AsyncClient(timeout=20) as client:
            response = await client.post(url, params={"key": api_key}, json=payload)
            response.raise_for_status()
            data = response.json()
        text = data["candidates"][0]["content"]["parts"][0]["text"]
        transcript = fallback_report.get("transcript_full") or (session.samples[-1].transcript if session.samples else "")
        report = merge_report_with_fallback(json.loads(text), fallback_report, transcript)
        report["used_gemini"] = True
        if isinstance(report.get("analysis_basis"), dict):
            report["analysis_basis"]["analysis_source"] = "gemini_plus_heuristic"
        report.setdefault("reference_speaker_comparison", fallback_report.get("reference_speaker_comparison"))
        return polish_user_report_text(sanitize_report_for_user(report))
    except Exception:
        fallback_report["summary"] = "기본 분석으로 리포트를 정리했습니다. 지금 연습에서 바로 고칠 부분을 중심으로 확인해 주세요."
        return polish_user_report_text(sanitize_report_for_user(fallback_report))


async def ask_gemini_for_report_v2(session: SessionState, fallback_report: dict[str, Any]) -> dict[str, Any]:
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    use_gemini_final_report = os.getenv("ENABLE_GEMINI_FINAL_REPORT", "false").lower() in {"1", "true", "yes"}
    if not api_key or not use_gemini_final_report:
        return polish_user_report_text(sanitize_report_for_user(fallback_report))

    prompt = build_ai_prompt_payload(session, fallback_report)
    payload = {
        "contents": [{"parts": [{"text": json.dumps(prompt, ensure_ascii=False)}]}],
        "generationConfig": {"responseMimeType": "application/json"},
    }

    try:
        data, remaining = await call_gemini_api(
            kind="final_report",
            payload=payload,
            timeout_seconds=6,
            count_against_limit=True,
            retry_on_unavailable=0,
        )
        text = (
            ((data or {}).get("candidates") or [{}])[0]
            .get("content", {})
            .get("parts", [{}])[0]
            .get("text")
        )
        if not text:
            return sanitize_report_for_user(fallback_report)
        transcript = fallback_report.get("transcript_full") or (session.samples[-1].transcript if session.samples else "")
        report = merge_report_with_fallback(json.loads(text), fallback_report, transcript)
        report["used_gemini"] = True
        if isinstance(report.get("analysis_basis"), dict):
            report["analysis_basis"]["analysis_source"] = "gemini_plus_heuristic"
        report["gemini_limits"] = {
            "rate_limited": False,
            "remaining_calls": remaining,
        }
        report.setdefault("reference_speaker_comparison", fallback_report.get("reference_speaker_comparison"))
        return polish_user_report_text(sanitize_report_for_user(report))
    except GeminiRateLimitError as exc:
        fallback_report["summary"] = "AI 분석이 잠시 막혀 기본 리포트로 정리했습니다. 지금 바로 고칠 부분부터 확인해 주세요."
        fallback_report["used_gemini"] = False
        fallback_report["gemini_limits"] = {
            "rate_limited": True,
            "retry_after_seconds": round(exc.retry_after, 1),
            "remaining_calls": exc.remaining,
        }
        return polish_user_report_text(sanitize_report_for_user(fallback_report))
    except Exception:
        fallback_report["summary"] = "기본 분석으로 리포트를 정리했습니다. 지금 연습에서 바로 고칠 부분을 중심으로 확인해 주세요."
        fallback_report["used_gemini"] = False
        return polish_user_report_text(sanitize_report_for_user(fallback_report))


@app.get("/api/ai/status")
async def ai_status(probe: bool = False) -> dict[str, Any]:
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    model = get_gemini_model()
    if not api_key:
        return {
            "configured": False,
            "live": False,
            "model": model,
            "message": "GEMINI_API_KEY가 비어 있어 로컬 분석만 사용합니다.",
        }

    if not probe:
        return {
            "configured": True,
            "live": gemini_runtime_state.get("last_live"),
            "model": model,
            "cached": False,
            "checked_live": False,
            "remaining_calls": gemini_limiter.remaining(),
            "last_ok_at": gemini_runtime_state.get("last_ok_at"),
            "last_error_at": gemini_runtime_state.get("last_error_at"),
            "last_error_code": gemini_runtime_state.get("last_error_code"),
            "last_error_message": gemini_runtime_state.get("last_error_message"),
            "last_call_kind": gemini_runtime_state.get("last_call_kind"),
            "message": "Gemini API 키가 설정되어 있습니다. 실제 외부 호출은 분석이 실행될 때만 진행됩니다.",
        }

    now = time.monotonic()
    cached = ai_status_cache.get("payload")
    cached_at = float(ai_status_cache.get("checked_at") or 0.0)
    if cached and now - cached_at < GEMINI_STATUS_CACHE_SECONDS:
        return {**cached, "cached": True}

    allowed, retry_after, remaining = gemini_limiter.allow()
    if not allowed:
        return {
            "configured": True,
            "live": False,
            "model": model,
            "cached": True,
            "rate_limited": True,
            "remaining_calls": remaining,
            "retry_after_seconds": round(retry_after, 1),
            "message": f"Gemini 상태 확인이 너무 자주 요청되어 캐시된 상태를 반환합니다. 약 {retry_after:.0f}초 후 다시 확인하세요.",
        }

    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    payload = {
        "contents": [{"parts": [{"text": "Return the single word OK."}]}],
        "generationConfig": {"maxOutputTokens": 8},
    }
    try:
        async with httpx.AsyncClient(timeout=12) as client:
            response = await client.post(url, params={"key": api_key}, json=payload)
            response.raise_for_status()
            data = response.json()
            text = data["candidates"][0]["content"]["parts"][0]["text"].strip()
            result = {
                "configured": True,
                "live": bool(text),
                "model": model,
                "rate_limited": False,
                "remaining_calls": remaining,
                "message": "Gemini API 연결이 정상입니다.",
            }
            ai_status_cache["checked_at"] = now
            ai_status_cache["payload"] = result
            return result
    except httpx.HTTPStatusError as exc:
        status_code = exc.response.status_code
        if status_code == 429:
            message = "Gemini API 할당량 또는 요청 제한으로 현재 AI 리포트가 대기 중입니다."
        elif status_code in {401, 403}:
            message = "Gemini API 키 권한을 확인해야 합니다."
        else:
            message = f"Gemini API 연결 확인에 실패했습니다. 상태 코드: {status_code}"
        result = {
            "configured": True,
            "live": False,
            "model": model,
            "rate_limited": False,
            "remaining_calls": remaining,
            "message": message,
        }
        ai_status_cache["checked_at"] = now
        ai_status_cache["payload"] = result
        return result
    except Exception:
        result = {
            "configured": True,
            "live": False,
            "model": model,
            "rate_limited": False,
            "remaining_calls": remaining,
            "message": "Gemini API 연결을 확인하지 못했습니다. 잠시 후 다시 시도하세요.",
        }
        ai_status_cache["checked_at"] = now
        ai_status_cache["payload"] = result
        return result


@app.get("/api/openai/status")
async def openai_status(probe: bool = False) -> dict[str, Any]:
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    model = OPENAI_AUDIENCE_MODEL
    if not api_key:
        return {
            "configured": False,
            "live": False,
            "model": model,
            "message": "OPENAI_API_KEY가 비어 있어 관객 채팅은 기본 문구로 대체됩니다.",
        }

    base_payload = {
        "configured": True,
        "live": openai_runtime_state.get("last_live"),
        "model": model,
        "remaining_calls": openai_audience_limiter.remaining(),
        "last_ok_at": openai_runtime_state.get("last_ok_at"),
        "last_error_at": openai_runtime_state.get("last_error_at"),
        "last_error_code": openai_runtime_state.get("last_error_code"),
        "last_error_message": openai_runtime_state.get("last_error_message"),
        "last_call_kind": openai_runtime_state.get("last_call_kind"),
    }
    if not probe:
        return {
            **base_payload,
            "checked_live": False,
            "message": "OpenAI API 키가 설정되어 있습니다. 실제 연결은 관객 채팅 또는 probe에서 확인합니다.",
        }

    payload = {
        "model": model,
        "input": "Return OK.",
        "max_output_tokens": 16,
    }
    try:
        async with httpx.AsyncClient(timeout=12) as client:
            response = await client.post(
                "https://api.openai.com/v1/responses",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload,
            )
            response.raise_for_status()
            text = extract_openai_text(response.json())
            mark_openai_success("status_probe")
            return {
                **base_payload,
                "live": bool(text),
                "checked_live": True,
                "last_ok_at": openai_runtime_state.get("last_ok_at"),
                "last_error_code": None,
                "last_error_message": None,
                "message": "OpenAI API 연결이 정상입니다.",
            }
    except httpx.HTTPStatusError as exc:
        status_code = exc.response.status_code
        detail = exc.response.text[:500]
        mark_openai_error("status_probe", detail or f"HTTP {status_code}", status_code)
        if status_code == 429:
            message = "OpenAI API 할당량 또는 요청 제한을 확인해야 합니다."
        elif status_code in {401, 403}:
            message = "OpenAI API 키 권한을 확인해야 합니다."
        elif status_code == 404:
            message = f"OpenAI 모델 '{model}'을 사용할 수 없습니다. 모델명을 확인해 주세요."
        else:
            message = f"OpenAI API 연결 확인에 실패했습니다. 상태 코드: {status_code}"
        return {
            **base_payload,
            "live": False,
            "checked_live": True,
            "last_error_code": status_code,
            "last_error_message": openai_runtime_state.get("last_error_message"),
            "message": message,
        }
    except Exception as exc:
        mark_openai_error("status_probe", str(exc))
        return {
            **base_payload,
            "live": False,
            "checked_live": True,
            "last_error_message": openai_runtime_state.get("last_error_message"),
            "message": "OpenAI API 연결을 확인하지 못했습니다. 네트워크 또는 키 설정을 확인해 주세요.",
        }

@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.post("/api/script/import", response_model=ImportedScriptResponse)
async def import_script_file(file: UploadFile = File(...)) -> ImportedScriptResponse:
    filename = file.filename or "script"
    content = await file.read()
    if len(content) > MAX_SCRIPT_FILE_BYTES:
        raise HTTPException(status_code=413, detail="10MB 이하의 파일만 불러올 수 있습니다.")

    text, source_type = extract_script_text(filename, content)
    normalized = normalize_imported_text(text)
    if not normalized:
        raise HTTPException(status_code=400, detail="파일에서 읽을 수 있는 대본 내용을 찾지 못했습니다.")

    return ImportedScriptResponse(
        filename=filename,
        text=normalized,
        character_count=len(normalized),
        source_type=source_type,
    )


@app.post("/api/reference/youtube")
async def preview_youtube_reference(payload: dict[str, str]) -> dict[str, Any]:
    reference = await build_reference_video(payload.get("url"))
    if not reference:
        raise HTTPException(status_code=400, detail="YouTube 영상 주소를 입력해 주세요.")
    return reference


async def collect_uploaded_materials(
    script: str,
    materials: list[UploadFile] | None,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    uploaded_blobs: list[dict[str, Any]] = []
    uploaded_materials: list[dict[str, Any]] = []
    for file in (materials or [])[:MAX_PRESENTATION_UPLOAD_FILES]:
        data = await file.read()
        filename = file.filename or "material"
        uploaded_blobs.append(
            {
                "filename": filename,
                "content_type": file.content_type,
                "data": data,
            }
        )
        uploaded_materials.append(analyze_material_file(script, filename, file.content_type, data))
    return uploaded_blobs, uploaded_materials


@app.post("/api/preflight")
async def preflight_session(
    script: str = Form(...),
    reference_video_url: str | None = Form(default=None),
    materials: list[UploadFile] | None = File(default=None),
) -> dict[str, Any]:
    normalized_script = script.strip()
    uploaded_blobs, uploaded_materials = await collect_uploaded_materials(normalized_script, materials)
    local_script_feedback = script_quality(normalized_script)
    preflight_feedback = await analyze_preflight_with_gemini(normalized_script, uploaded_materials, uploaded_blobs)
    uploaded_materials = merge_vision_material_feedback(
        uploaded_materials,
        (preflight_feedback or {}).get("presentation_material"),
    )
    material_feedback = merge_material_feedback_summary(
        build_material_feedback(normalized_script, uploaded_materials),
        (preflight_feedback or {}).get("presentation_material"),
    )
    reference_video = await build_reference_video(reference_video_url)
    return {
        "script_feedback": merge_gemini_script_feedback(
            local_script_feedback,
            (preflight_feedback or {}).get("script_feedback"),
        ),
        "criteria_basis": presentation_criteria(),
        "presentation_material": material_feedback,
        "reference_video": reference_video,
        "gemini_ready": bool(os.getenv("GEMINI_API_KEY", "").strip()),
        "upload_count": len(uploaded_materials),
    }


@app.post("/api/session/start")
async def start_session(
    script: str = Form(...),
    reference_video_url: str | None = Form(default=None),
    materials: list[UploadFile] | None = File(default=None),
) -> dict[str, Any]:
    session_id = str(uuid4())
    normalized_script = script.strip()
    _uploaded_blobs, uploaded_materials = await collect_uploaded_materials(normalized_script, materials)
    reference_video = await build_reference_video(reference_video_url)

    session = SessionState(
        id=session_id,
        script=normalized_script,
        created_at=datetime.now(timezone.utc).isoformat(),
        materials=uploaded_materials,
        reference_video=reference_video,
    )
    sessions[session_id] = session
    save_session(session)
    return {
        "session_id": session_id,
        "criteria_basis": presentation_criteria(),
        "reference_video": reference_video,
        "gemini_ready": bool(os.getenv("GEMINI_API_KEY", "").strip()),
        "upload_count": len(uploaded_materials),
    }

@app.post("/api/session/{session_id}/metric")
def add_metric(session_id: str, sample: MetricSample) -> dict[str, str]:
    session = get_session_state(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    session.samples.append(sample)
    persist_metric(session_id, sample)
    return {"status": "stored"}


@app.post("/api/session/{session_id}/audience/chat")
async def create_audience_chat(session_id: str, request: AudienceChatRequest) -> dict[str, Any]:
    session = get_session_state(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    return await build_audience_chat(session_id, session, request)


@app.post("/api/session/{session_id}/finish")
async def finish_session(session_id: str, request: FinishSessionRequest) -> dict[str, Any]:
    session = get_session_state(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    transcript = request.transcript or (session.samples[-1].transcript if session.samples else "")
    fallback = build_heuristic_report(session, transcript)
    return await ask_gemini_for_report_v2(session, fallback)
