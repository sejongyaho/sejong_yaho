import os
import re
from io import BytesIO
from typing import Any

from docx import Document
from fastapi import HTTPException
from pptx import Presentation
from pypdf import PdfReader

from ..config import SUPPORTED_SCRIPT_EXTENSIONS, TEXT_EXTENSIONS


def tokenize(text: str) -> list[str]:
    return re.findall(r"[가-힣A-Za-z0-9']+", text.lower())


def count_syllables(text: str) -> int:
    hangul = len(re.findall(r"[가-힣]", text))
    latin_words = re.findall(r"[A-Za-z0-9']+", text)
    return hangul + sum(max(1, round(len(word) / 3)) for word in latin_words)


def clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, value))


def score_distance(value: float, target: float, tolerance: float, floor: float = 35) -> float:
    if value <= 0:
        return floor
    return clamp(100 - abs(value - target) / tolerance * 35, floor, 100)


def format_seconds(seconds: float) -> str:
    whole = max(0, int(seconds))
    minutes = whole // 60
    rest = whole % 60
    return f"{minutes:02d}:{rest:02d}"


def transcript_delta(previous: str, current: str) -> str:
    previous = previous.strip()
    current = current.strip()
    if not current:
        return ""
    if previous and current.startswith(previous):
        return current[len(previous) :].strip()
    previous_tokens = tokenize(previous)
    current_tokens = tokenize(current)
    if not previous_tokens:
        return current
    shared = 0
    for index, token in enumerate(current_tokens):
        if index < len(previous_tokens) and previous_tokens[index] == token:
            shared += 1
        else:
            break
    if shared >= len(previous_tokens) - 2:
        return " ".join(current_tokens[shared:]).strip()
    return current[-180:].strip()


def recent_excerpt(text: str, limit: int = 140) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    if len(compact) <= limit:
        return compact
    return f"...{compact[-limit:]}"


def normalize_imported_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def file_extension(filename: str) -> str:
    _, extension = os.path.splitext(filename.lower())
    return extension


def extract_text_from_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(f"[{index}페이지]\n{page_text.strip()}")
    return "\n\n".join(pages)


def extract_text_from_docx(content: bytes) -> str:
    document = Document(BytesIO(content))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_pptx(content: bytes) -> str:
    deck = Presentation(BytesIO(content))
    slides = []
    for index, slide in enumerate(deck.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            slides.append(f"[슬라이드 {index}]\n" + "\n".join(slide_text))
    return "\n\n".join(slides)


def extract_script_text(filename: str, content: bytes) -> tuple[str, str]:
    extension = file_extension(filename)
    if extension not in SUPPORTED_SCRIPT_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_SCRIPT_EXTENSIONS))
        raise HTTPException(status_code=400, detail=f"지원하지 않는 파일 형식입니다. 지원 형식: {supported}")

    try:
        if extension in TEXT_EXTENSIONS:
            try:
                return content.decode("utf-8-sig"), "text"
            except UnicodeDecodeError:
                return content.decode("cp949", errors="ignore"), "text"
        if extension == ".pdf":
            return extract_text_from_pdf(content), "pdf"
        if extension == ".docx":
            return extract_text_from_docx(content), "docx"
        if extension == ".pptx":
            return extract_text_from_pptx(content), "pptx"
    except Exception as exc:
        raise HTTPException(status_code=400, detail="파일에서 대본 텍스트를 읽지 못했습니다.") from exc

    raise HTTPException(status_code=400, detail="파일을 읽지 못했습니다.")


def script_quality(script: str) -> dict[str, Any]:
    words = tokenize(script)
    sentence_count = max(1, len(re.findall(r"[.!?]+", script)))
    avg_sentence_words = len(words) / sentence_count if words else 0
    unique_ratio = len(set(words)) / len(words) if words else 0
    has_opening = any(token in script for token in ["안녕하세요", "오늘", "소개", "주제"])
    has_closing = any(token in script for token in ["감사합니다", "정리", "결론", "마치"])

    score = 55
    score += 15 if 8 <= avg_sentence_words <= 22 else 4
    score += 10 if unique_ratio > 0.45 else 4
    score += 8 if has_opening else 0
    score += 8 if has_closing else 0
    score += 4 if len(words) >= 80 else 0

    suggestions = []
    if avg_sentence_words > 24:
        suggestions.append("긴 문장을 둘로 나누면 청중이 전달을 따라가기 쉬워집니다.")
    if not has_opening:
        suggestions.append("처음 15초 안에 발표 주제와 듣는 이유를 분명하게 말해보세요.")
    if not has_closing:
        suggestions.append("마지막에 핵심 요약과 감사 인사를 짧게 넣어 마무리감을 주세요.")
    if len(words) < 80:
        suggestions.append("대본이 짧아 리허설 피드백이 제한될 수 있어요. 전환 문장을 조금 더 넣어보세요.")

    summary = "발표 대본의 구조와 문장 길이, 시작과 마무리 문장을 기준으로 점검했습니다."
    if 8 <= avg_sentence_words <= 22 and has_opening and has_closing:
        summary = "대본 구조가 비교적 안정적이고, 발표 시작과 마무리 흐름도 잘 갖춰져 있습니다."
    elif avg_sentence_words > 24:
        summary = "문장이 길어져 전달 밀도가 높습니다. 끊어 말할 지점을 조금 더 만들어 주면 좋습니다."

    return {
        "score": round(clamp(score, 0, 100)),
        "word_count": len(words),
        "average_sentence_words": round(avg_sentence_words, 1),
        "summary": summary,
        "suggestions": suggestions[:4],
    }
